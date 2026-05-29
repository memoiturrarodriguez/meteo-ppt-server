# -*- coding: utf-8 -*-
"""
Generador PPT - Briefing meteorologico operacional estilo boletin tactico.
Compatible con JSON exportado desde Meteo Operaciones Especiales.

Reemplaza el archivo generar_ppt_meteo_final_v3.py si quieres usarlo como version principal.
Salida: METEO_<SECTOR>_<PERIODO>.pptx
"""

import json
import math
from pathlib import Path
from collections import defaultdict
from datetime import datetime, timedelta

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from PIL import Image, ImageDraw, ImageFont
except Exception:
    Image = None
    ImageDraw = None
    ImageFont = None


# ============================================================
# CONFIGURACION
# ============================================================
BASE_DIR = Path(__file__).parent
JSON_FILE = BASE_DIR / "reporte_meteo.json"
OUTPUT_FILE = BASE_DIR / "briefing_meteo_operacional_final_v3.pptx"

# ============================================================
# COLORES
# ============================================================
BG = RGBColor(9, 15, 24)
BG2 = RGBColor(13, 22, 34)
PANEL = RGBColor(22, 33, 48)
PANEL2 = RGBColor(29, 43, 61)
PANEL3 = RGBColor(36, 52, 72)
DARK = RGBColor(7, 11, 17)
WHITE = RGBColor(245, 248, 252)
MUTED = RGBColor(176, 190, 205)
BLUE = RGBColor(83, 169, 255)
CYAN = RGBColor(79, 218, 255)
GREEN = RGBColor(35, 176, 95)
GREEN2 = RGBColor(21, 122, 74)
YELLOW = RGBColor(246, 196, 65)
ORANGE = RGBColor(236, 144, 42)
RED = RGBColor(214, 54, 66)
RED2 = RGBColor(155, 35, 47)
BLACK = RGBColor(0, 0, 0)
PURPLE = RGBColor(123, 92, 255)

STATUS_COLORS = {
    "OPTIMO": RGBColor(38, 202, 126),
    "FAVORABLE": GREEN,
    "FAV": GREEN,
    "GO": GREEN,
    "MARGINAL": ORANGE,
    "MARG": ORANGE,
    "RESTRINGIDO": RED,
    "REST": RED,
    "NO GO": RED,
    "SIN EVALUAR": MUTED,
}

# ============================================================
# UTILIDADES GENERALES
# ============================================================

def load_json():
    with open(JSON_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def normalize(text):
    return (
        str(text or "")
        .upper()
        .replace("Á", "A").replace("É", "E").replace("Í", "I")
        .replace("Ó", "O").replace("Ú", "U").replace("Ñ", "N")
    )


def simple(text):
    return (
        str(text or "")
        .lower()
        .replace("á", "a").replace("é", "e").replace("í", "i")
        .replace("ó", "o").replace("ú", "u").replace("ñ", "n")
    )


def valnum(v, default=None):
    if isinstance(v, (int, float)) and not isinstance(v, bool):
        return float(v)
    if v is None:
        return default
    txt = str(v).replace(",", ".")
    out = ""
    dot = False
    sign = False
    for ch in txt:
        if ch.isdigit():
            out += ch
        elif ch == "." and not dot:
            out += ch
            dot = True
        elif ch == "-" and not sign and not out:
            out += ch
            sign = True
        elif out:
            break
    try:
        return float(out) if out not in ("", "-", ".") else default
    except Exception:
        return default


def fmt_num(v, dec=0, suffix=""):
    n = valnum(v, None)
    if n is None:
        return "-"
    if dec == 0:
        return f"{round(n):.0f}{suffix}"
    return f"{n:.{dec}f}{suffix}"


def format_dt(value, mode="time"):
    try:
        txt = str(value or "").replace("Z", "")
        dt = datetime.fromisoformat(txt)

        meses = [
            "JAN", "FEB", "MAR", "APR", "MAY", "JUN",
            "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"
        ]

        if mode in ("date", "daylabel", "filename"):
            return f"{dt.day:02d}{meses[dt.month-1]}{dt.year}"

        return dt.strftime("%H")
    except Exception:
        txt = str(value or "")
        try:
            raw = txt[:10]
            dt = datetime.fromisoformat(raw)
            meses = [
                "JAN", "FEB", "MAR", "APR", "MAY", "JUN",
                "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"
            ]
            if mode in ("date", "daylabel", "filename"):
                return f"{dt.day:02d}{meses[dt.month-1]}{dt.year}"
        except Exception:
            pass

        if mode in ("date", "daylabel", "filename"):
            return txt[:10]
        return txt[11:13] if len(txt) >= 13 else txt

def grouped_by_date(forecast):
    grouped = defaultdict(list)
    for b in forecast or []:
        grouped[format_dt(b.get("hora"), "date")].append(b)
    return dict(grouped)


def periodo_label(data, forecast):
    periodo = data.get("periodo_briefing") or {}
    ini = periodo.get("inicio")
    fin = periodo.get("termino")
    if ini and fin:
        return f"{format_dt(ini, 'date')} a {format_dt(fin, 'date')}"
    if forecast:
        return f"{format_dt(forecast[0].get('hora'), 'date')} a {format_dt(forecast[-1].get('hora'), 'date')}"
    return "Periodo no informado"

def periodo_tipo(data):
    periodo = data.get("periodo_briefing") or {}
    if periodo.get("mayor_72h"):
        return "PROYECCIÓN PRELIMINAR PARA PLANIFICACIÓN"
    return periodo.get("tipo") or "PROYECCIÓN TÁCTICA"

def _parse_period_date(value, end_of_day=False):
    """Acepta ISO, yyyy-mm-dd, dd/mm/yyyy, dd-mm-yyyy; retorna datetime o None."""
    if not value:
        return None
    if isinstance(value, datetime):
        dt = value
    else:
        txt = str(value).strip().replace("Z", "")
        dt = None
        # ISO primero
        try:
            dt = datetime.fromisoformat(txt)
        except Exception:
            pass
        if dt is None:
            for fmt in ("%d/%m/%Y", "%d-%m-%Y", "%Y/%m/%d", "%Y-%m-%d"):
                try:
                    dt = datetime.strptime(txt[:10], fmt)
                    break
                except Exception:
                    continue
        if dt is None:
            return None
    if end_of_day:
        return dt.replace(hour=23, minute=59, second=59, microsecond=999999)
    return dt.replace(hour=0, minute=0, second=0, microsecond=0)


def filtered_forecast(data):
    """Respeta fecha inicio/término del briefing y mantiene solo el período solicitado."""
    forecast = list(data.get("forecast72h", []) or [])
    periodo = data.get("periodo_briefing") or {}
    inicio = _parse_period_date(periodo.get("inicio") or periodo.get("fecha_inicio") or periodo.get("desde"), False)
    termino = _parse_period_date(periodo.get("termino") or periodo.get("fecha_termino") or periodo.get("hasta"), True)
    if not inicio and not termino:
        return forecast
    out = []
    for b in forecast:
        dt = _parse_dt_safe(b.get("hora")) or _parse_period_date(b.get("hora"))
        if not dt:
            out.append(b)
            continue
        if inicio and dt < inicio:
            continue
        if termino and dt > termino:
            continue
        out.append(b)
    return out


def _period_bounds(data):
    periodo = data.get("periodo_briefing") or {}
    inicio = _parse_period_date(periodo.get("inicio") or periodo.get("fecha_inicio") or periodo.get("desde"), False)
    termino = _parse_period_date(periodo.get("termino") or periodo.get("fecha_termino") or periodo.get("hasta"), True)
    return inicio, termino


def _placeholder_block(dt):
    """Bloque visible si el JSON no trae datos para una fecha solicitada.
    Evita que el PPT invente una condición favorable cuando faltan datos.
    """
    return {
        "hora": dt.isoformat(),
        "temperatura": "-",
        "viento": "-",
        "direccion_grados": "-",
        "direccion": "-",
        "rachas": "-",
        "nubosidad": "-",
        "precipitacion": "-",
        "visibilidad": "-",
        "base_nubes_ft": "-",
        "techo": "-",
        "isoterma_0_m": "-",
        "nieve_cm_h": "-",
        "nieve_acumulada_cm": "-",
        "weather_code": -1,
        "estado_general": "SIN DATOS",
        "estado": "SIN DATOS",
        "operaciones": [],
        "montana": {"estado": "SIN DATOS", "motivo": "Sin datos para este bloque"},
        "sin_datos": True,
    }


def force_exact_requested_period(data):
    """Hace que la presentación use exactamente las fechas elegidas en la app.
    Ejemplo: si el operador selecciona 22-25 MAY, se eliminan bloques del 21
    y se conserva/crea la grilla 22, 23, 24 y 25 con horarios 02-23.
    """
    data = dict(data or {})
    inicio, termino = _period_bounds(data)
    forecast = filtered_forecast(data)

    if not inicio or not termino:
        data["forecast72h"] = forecast
        return data

    hours_allowed = [2, 5, 8, 11, 14, 17, 20, 23]
    existing = {}
    for b in forecast:
        dt = _parse_dt_safe(b.get("hora")) or _parse_period_date(b.get("hora"))
        if not dt:
            continue
        key = (dt.year, dt.month, dt.day, dt.hour)
        existing[key] = b

    exact = []
    day = datetime(inicio.year, inicio.month, inicio.day)
    last_day = datetime(termino.year, termino.month, termino.day)
    while day <= last_day:
        for h in hours_allowed:
            dt = datetime(day.year, day.month, day.day, h, 0, 0)
            key = (dt.year, dt.month, dt.day, dt.hour)
            exact.append(existing.get(key, _placeholder_block(dt)))
        day = day + timedelta(days=1)

    data["forecast72h"] = exact
    # Etiqueta útil para portada/conclusiones.
    data.setdefault("periodo_briefing", {})["horas_presentacion"] = len(exact) * 3
    return data


def apply_period_filter(data):
    """Mutación controlada para que todas las slides usen exactamente los días solicitados."""
    return force_exact_requested_period(data)


def period_hours_from_forecast(forecast):
    dts = [_parse_dt_safe(b.get("hora")) for b in forecast or []]
    dts = [d for d in dts if d]
    if not dts:
        return 0
    # Bloques de 3h; sumar 3h para incluir el último bloque.
    return int(round((max(dts) - min(dts)).total_seconds() / 3600.0)) + 3


def status_color(estado):
    e = normalize(estado)
    if "OPTIMO" in e:
        return STATUS_COLORS["OPTIMO"]
    if e == "GO" or "FAVORABLE" in e or e == "FAV":
        return GREEN
    if "MARGINAL" in e or "RESTRICC" in e and "CON" in e or e == "MARG":
        return ORANGE
    if "RESTRINGIDO" in e or "NO GO" in e or e == "REST":
        return RED
    return MUTED



# ============================================================
# REGLAS OPERACIONALES FORZADAS PARA PPT
# ============================================================

def is_aerial_operation(opname):
    n = normalize(opname)
    return (
        "SALTO BASICO" in n or
        "LANZAMIENTO" in n or
        "SALTO LIBRE" in n or
        "EVACAM" in n or
        n in ("SBM", "LC", "SLM")
    )


def _op_thresholds(opname):
    """Límites operacionales en kt. Se evalúa viento sostenido y racha."""
    n = normalize(opname)
    if "SALTO BASICO" in n or n == "SBM":
        return {"restrict": 13, "marginal": 12}
    if "LANZAMIENTO" in n or n == "LC":
        return {"restrict": 15, "marginal": 14}
    if "SALTO LIBRE" in n or n == "SLM":
        return {"restrict": 17, "marginal": 16}
    if "EVACAM" in n:
        return {"restrict": 25, "marginal": 15}
    return {"restrict": None, "marginal": None}


def _state_rank(estado):
    e = normalize(estado)
    if "RESTRINGIDO" in e or "NO GO" in e or e == "REST":
        return 2
    if "MARGINAL" in e or e == "MARG":
        return 1
    return 0


def _rank_state(rank):
    if rank >= 2:
        return "RESTRINGIDO"
    if rank == 1:
        return "MARGINAL"
    return "FAVORABLE"


def _aerial_forced_rank(block, opname):
    """
    Regla táctica para la presentación:
    - Si la visibilidad es crítica, TODAS las operaciones aéreas quedan restringidas.
    - Si la visibilidad es reducida, TODAS las operaciones aéreas quedan marginales.
    - Se suman restricciones por precipitación, techo bajo y límites de viento/rachas.
    """
    if not is_aerial_operation(opname):
        return 0

    rank = 0
    vis = valnum(block.get("visibilidad"), None)
    viento = valnum(block.get("viento"), 0) or 0
    rachas = valnum(block.get("rachas"), 0) or 0
    precip = valnum(block.get("precipitacion"), 0) or 0
    techo = valnum(get_metric(block, ["techo", "cloud_base", "base_nubes", "baseNubes"], None), None)

    # Visibilidad: criterio común para toda operación aérea.
    # Si el valor viene como 0 por redondeo, se considera crítico.
    if vis is not None:
        if vis <= 5:
            rank = max(rank, 2)
        elif vis <= 10:
            rank = max(rank, 1)

    # Techo bajo crítico para aeronaves.
    if techo is not None:
        if techo > 0 and techo < 500:
            rank = max(rank, 2)
        elif 500 <= techo < 1000:
            rank = max(rank, 1)

    # Precipitación: compatible con reglas usadas por la app.
    if precip >= 1.0:
        rank = max(rank, 2)
    elif precip >= 0.1:
        rank = max(rank, 1)

    # Viento/rachas por tipo de operación.
    th = _op_thresholds(opname)
    restrict = th.get("restrict")
    marginal = th.get("marginal")
    wind_ref = max(viento, rachas)
    if restrict is not None and wind_ref >= restrict:
        rank = max(rank, 2)
    elif marginal is not None and wind_ref >= marginal:
        rank = max(rank, 1)

    return rank


def _non_aerial_forced_rank(block, opname):
    """Reglas mínimas para operaciones no aéreas sin sobre-restringir."""
    if isinstance(block, dict) and block.get("sin_datos"):
        return 0
    n = normalize(opname)
    vis = valnum(block.get("visibilidad"), None)
    temp = valnum(block.get("temperatura"), None)
    rank = 0
    if "MARCHA" in n:
        if vis is not None and vis <= 5:
            rank = max(rank, 1)
        if temp is not None and temp >= 25:
            rank = max(rank, 1)
    return rank


def forced_operation_state(block, opname, original_state="FAVORABLE"):
    if isinstance(block, dict) and block.get("sin_datos"):
        return "SIN DATOS"
    base_rank = _state_rank(original_state)
    forced_rank = max(
        _aerial_forced_rank(block, opname),
        _non_aerial_forced_rank(block, opname),
    )
    return _rank_state(max(base_rank, forced_rank))


def forced_operation_reason(block, opname):
    """Motivo compacto para tarjetas/leyendas si el JSON no trae el motivo correcto."""
    motivos = []
    n = normalize(opname)
    vis = valnum(block.get("visibilidad"), None)
    viento = valnum(block.get("viento"), 0) or 0
    rachas = valnum(block.get("rachas"), 0) or 0
    precip = valnum(block.get("precipitacion"), 0) or 0
    techo = valnum(get_metric(block, ["techo", "cloud_base", "base_nubes", "baseNubes"], None), None)

    if is_aerial_operation(opname):
        if vis is not None and vis <= 5:
            motivos.append(f"Visibilidad crítica ({vis:g} km)")
        elif vis is not None and vis <= 10:
            motivos.append(f"Visibilidad reducida ({vis:g} km)")

        if techo is not None and techo > 0 and techo < 500:
            motivos.append(f"Techo/base de nubes <500 ft ({techo:g} ft)")
        elif techo is not None and 500 <= techo < 1000:
            motivos.append(f"Techo/base de nubes marginal ({techo:g} ft)")

        if precip >= 1.0:
            motivos.append(f"Precipitación >=1.0 mm ({precip:g} mm)")
        elif precip >= 0.1:
            motivos.append(f"Precipitación marginal ({precip:g} mm)")

        th = _op_thresholds(opname)
        restrict = th.get("restrict")
        marginal = th.get("marginal")
        wind_ref = max(viento, rachas)
        if restrict is not None and wind_ref >= restrict:
            motivos.append(f"Viento/rachas sobre límite {restrict} kt ({wind_ref:g} kt)")
        elif marginal is not None and wind_ref >= marginal:
            motivos.append(f"Viento/rachas marginales ({wind_ref:g} kt)")

    elif "MARCHA" in n:
        if vis is not None and vis <= 5:
            motivos.append(f"Visibilidad reducida ({vis:g} km)")
        temp = valnum(block.get("temperatura"), None)
        if temp is not None and temp >= 25:
            motivos.append(f"Temperatura >=25°C ({temp:g}°C)")

    return " / ".join(motivos)

def estado_general_bloque(b):
    if isinstance(b, dict) and b.get("sin_datos"):
        return "SIN DATOS"

    # La condición general de la presentación debe respetar restricciones comunes
    # de seguridad aeronáutica, incluso si el JSON de la app viene incompleto.
    ops_ref = ["SALTO BASICO", "LANZAMIENTO", "SALTO LIBRE", "EVACAM", "MARCHA"]
    ranks = []
    for op in ops_ref:
        ranks.append(_state_rank(get_op_state(b, op)))
    if ranks:
        return _rank_state(max(ranks))

    estado = b.get("estado_general") or b.get("estado") or ""
    if estado:
        e = normalize(estado)
        if e == "GO":
            return "FAVORABLE"
        if "NO GO" in e:
            return "RESTRINGIDO"
        return estado

    return peor_estado_ops(b.get("operaciones", []))


def peor_estado_ops(ops):
    worst = "FAVORABLE"
    for op in ops or []:
        e = normalize(op.get("estado", ""))
        if "RESTRINGIDO" in e:
            return "RESTRINGIDO"
        if "MARGINAL" in e:
            worst = "MARGINAL"
    return worst


def estado_abbr(estado):
    e = normalize(estado)
    if "RESTRINGIDO" in e or "NO GO" in e:
        return "REST"
    if "MARGINAL" in e:
        return "MARG"
    if "FAVORABLE" in e or e == "GO":
        return "FAV"
    return "-"


def op_abbr(name):
    n = normalize(name)
    if "SALTO BASICO" in n:
        return "SBM"
    if "LANZAMIENTO" in n:
        return "LC"
    if "SALTO LIBRE" in n:
        return "SLM"
    if "EVACAM" in n:
        return "EVACAM"
    if "MARCHA" in n:
        return "MARCHA"
    if "EMBARC" in n:
        return "EMB"
    if "MONT" in n:
        return "MONTAÑA"
    return str(name or "-")[:8].upper()


def get_humedad(b):
    for k in ("humedad", "humedad_relativa", "humedadRelativa", "hr", "HR"):
        if k in b:
            return valnum(b.get(k), None)
    return None


def get_op_state(block, opname):
    if isinstance(block, dict) and block.get("sin_datos"):
        return "SIN DATOS"
    target = normalize(opname)
    original = "FAVORABLE"
    for op in block.get("operaciones", []) or []:
        if target in normalize(op.get("operacion", "")):
            original = op.get("estado", "FAVORABLE")
            break
    if target == "MONTAÑA" and isinstance(block.get("montana"), dict):
        original = block["montana"].get("estado", original)
    return forced_operation_state(block, opname, original)


def get_metric(block, keys, default="-"):
    for k in keys:
        if k in block and block.get(k) is not None:
            return block.get(k)
    return default


def text_wrap(s, max_chars=130):
    s = str(s or "")
    if len(s) <= max_chars:
        return s
    return s[: max_chars - 3].rstrip() + "..."


def _safe_filename_part(value, fallback="SECTOR"):
    """Convierte texto libre del sector a un nombre seguro de archivo."""
    raw = str(value or fallback).strip()
    if not raw:
        raw = fallback

    # Si viene "Peldehue, Región Metropolitana..." usar solo el primer tramo.
    raw = raw.split(",")[0].strip() or fallback

    replacements = {
        "á": "a", "é": "e", "í": "i", "ó": "o", "ú": "u", "ü": "u", "ñ": "n",
        "Á": "A", "É": "E", "Í": "I", "Ó": "O", "Ú": "U", "Ü": "U", "Ñ": "N",
    }
    for a, b in replacements.items():
        raw = raw.replace(a, b)

    safe = []
    last_us = False
    for ch in raw:
        if ch.isalnum():
            safe.append(ch.upper())
            last_us = False
        elif ch in (" ", "-", "_", "."):
            if not last_us:
                safe.append("_")
                last_us = True
        # cualquier otro símbolo se omite

    out = "".join(safe).strip("_")
    return out[:45] if out else fallback


def _period_for_filename(data):
    """Devuelve rango 22MAY2026_25MAY2026 desde periodo_briefing o forecast."""
    periodo = data.get("periodo_briefing") or {}
    forecast = data.get("forecast72h", []) or []

    ini = periodo.get("inicio")
    fin = periodo.get("termino")

    if not ini and forecast:
        ini = forecast[0].get("hora")
    if not fin and forecast:
        fin = forecast[-1].get("hora")

    a = format_dt(ini, "filename") if ini else ""
    b = format_dt(fin, "filename") if fin else ""

    if a and b:
        return f"{a}_{b}"
    if a:
        return a
    return datetime.now().strftime("%d%b%Y_%H%M").upper()

def output_file_for_data(data):
    """Nombre automático: METEO_<SECTOR>_<PERIODO>.pptx"""
    sector = (
        data.get("sector")
        or data.get("sector_mostrado")
        or data.get("nombre_sector")
        or data.get("ubicacion")
        or "SECTOR"
    )
    sector_safe = _safe_filename_part(sector)
    periodo_safe = _period_for_filename(data)
    return BASE_DIR / f"METEO_{sector_safe}_{periodo_safe}.pptx"


# ============================================================
# LUNA / NOCTURNIDAD
# ============================================================

def parse_date_any(value):
    """Convierte fecha/hora ISO o yyyy-mm-dd a datetime. Fallback: hoy."""
    try:
        txt = str(value or "").replace("Z", "")
        if len(txt) == 10:
            return datetime.fromisoformat(txt + "T00:00:00")
        return datetime.fromisoformat(txt)
    except Exception:
        return datetime.now()


def moon_info_for_date(dt):
    """
    Cálculo lunar aproximado sin depender de librerías externas.
    Suficiente para briefing operacional preliminar: fase, edad e iluminación estimada.
    """
    known_new_moon = datetime(2000, 1, 6, 18, 14)  # luna nueva referencia
    synodic_month = 29.53058867
    days = (dt - known_new_moon).total_seconds() / 86400.0
    age = days % synodic_month
    phase_angle = 2 * math.pi * age / synodic_month
    illumination = (1 - math.cos(phase_angle)) / 2 * 100

    if age < 1.8 or age >= 27.7:
        phase_name = "Luna nueva"
    elif age < 6.4:
        phase_name = "Creciente inicial"
    elif age < 8.9:
        phase_name = "Cuarto creciente"
    elif age < 13.8:
        phase_name = "Gibosa creciente"
    elif age < 16.8:
        phase_name = "Luna llena"
    elif age < 21.1:
        phase_name = "Gibosa menguante"
    elif age < 23.6:
        phase_name = "Cuarto menguante"
    else:
        phase_name = "Menguante final"

    return {
        "fase": phase_name,
        "iluminacion": round(illumination, 1),
        "edad_dias": round(age, 1),
    }


def lunar_eval_from_illum(illum):
    illum = valnum(illum, None)
    if illum is None:
        return "NO EVALUADA", "Sin dato de iluminación lunar. Mantener control de crepúsculos y observación directa."
    if illum <= 20:
        return "ÓPTIMA PARA INFIL", "Baja iluminación lunar: favorece ocultamiento, infiltración y movimientos discretos nocturnos."
    if illum <= 50:
        return "FAVORABLE", "Iluminación lunar moderada-baja: permite infiltración con buen control de exposición visual."
    if illum <= 75:
        return "MARGINAL", "Iluminación lunar relevante: aumenta exposición visual; evitar siluetas y zonas abiertas."
    return "DESFAVORABLE", "Alta iluminación lunar: desfavorable para infiltración discreta; priorizar cobertura, sombras y rutas ocultas."


def get_lunar_source(data):
    """Lee luna desde JSON si existe; si no, calcula desde el período del briefing."""
    lunar = data.get("luna") or data.get("moon") or data.get("nocturnidad") or {}
    periodo = data.get("periodo_briefing") or {}
    forecast = data.get("forecast72h", []) or []

    # 1) Formato nuevo del Dart: luna.resumen / luna.periodo
    if isinstance(lunar, dict) and lunar:
        resumen = lunar.get("resumen") or lunar.get("actual") or lunar
        illum = resumen.get("iluminacion") or resumen.get("iluminacion_promedio") or resumen.get("moon_illumination")
        fase = resumen.get("fase") or resumen.get("phase") or resumen.get("fase_lunar")
        periodo_lunar = lunar.get("periodo") if isinstance(lunar.get("periodo"), list) else []
        if illum is not None or fase or periodo_lunar:
            if illum is None and periodo_lunar:
                vals = [valnum(x.get("iluminacion"), None) for x in periodo_lunar if isinstance(x, dict)]
                vals = [v for v in vals if v is not None]
                illum = sum(vals) / len(vals) if vals else None
            if not fase and periodo_lunar:
                fase = periodo_lunar[0].get("fase")
            return {
                "fase": fase or "Fase no informada",
                "iluminacion": round(valnum(illum, 0), 1) if illum is not None else None,
                "periodo": periodo_lunar,
                "origen": "json",
            }

    # 2) Si el forecast trae campos lunares por bloque
    vals = []
    fases = []
    for b in forecast:
        illum = b.get("luna_iluminacion") or b.get("moon_illumination") or b.get("iluminacion_lunar")
        if illum is not None:
            v = valnum(illum, None)
            if v is not None:
                vals.append(v)
        fase = b.get("luna_fase") or b.get("moon_phase") or b.get("fase_lunar")
        if fase:
            fases.append(fase)
    if vals:
        return {
            "fase": fases[0] if fases else "Fase lunar estimada",
            "iluminacion": round(sum(vals) / len(vals), 1),
            "periodo": [],
            "origen": "forecast",
        }

    # 3) Fallback: cálculo aproximado por fecha inicio
    start = periodo.get("inicio") or (forecast[0].get("hora") if forecast else None) or datetime.now().isoformat()
    info = moon_info_for_date(parse_date_any(start))
    info["periodo"] = []
    info["origen"] = "calculado"
    return info


def build_lunar_text(data):
    info = get_lunar_source(data)
    illum = info.get("iluminacion")
    fase = info.get("fase") or "Fase lunar estimada"
    estado, evaluacion = lunar_eval_from_illum(illum)
    illum_txt = f"{illum:.1f}%" if isinstance(illum, (int, float)) else "N/D"

    extra = ""
    if info.get("origen") == "calculado":
        extra = "\nDato calculado de forma aproximada según fecha del briefing."

    return (
        f"Fase: {fase}\n"
        f"Iluminación estimada: {illum_txt}\n"
        f"Condición nocturna: {estado}\n"
        f"{evaluacion}"
        f"{extra}"
    )

# ============================================================
# PPT HELPERS
# ============================================================

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_rect(slide, x, y, w, h, fill=PANEL, line=RGBColor(60, 76, 96), radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.75)
    return shp


def add_text(slide, text, x, y, w, h, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = str(line)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, data, title, subtitle=""):
    sector = str(data.get("sector", "SECTOR") or "SECTOR").upper()
    coords = data.get("coordenadas", "")
    forecast = data.get("forecast72h", [])
    add_text(slide, "BOLETÍN METEOROLÓGICO", 0.35, 0.18, 4.0, 0.25, 11, True, CYAN)
    add_text(slide, f"{sector} · {periodo_label(data, forecast)}", 0.35, 0.42, 7.0, 0.22, 9.5, False, MUTED)
    add_text(slide, title, 0.35, 0.72, 8.5, 0.38, 19, True, WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.36, 1.07, 8.5, 0.22, 9, False, MUTED)
    add_text(slide, "M · METEO | BOE", 10.65, 0.20, 2.3, 0.25, 10, True, WHITE, PP_ALIGN.RIGHT)
    add_text(slide, "Pelotón de Inteligencia", 10.65, 0.43, 2.3, 0.22, 8.5, False, MUTED, PP_ALIGN.RIGHT)
    add_text(slide, str(coords), 8.7, 1.02, 4.25, 0.22, 8.5, False, MUTED, PP_ALIGN.RIGHT)
    add_rect(slide, 0.35, 1.32, 12.65, 0.02, BLUE, BLUE, False)


def add_badge(slide, text, x, y, w, h, color):
    shp = add_rect(slide, x, y, w, h, color, color)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return shp


def set_cell(cell, text, size=7, bold=False, fill=None, color=WHITE, align=PP_ALIGN.CENTER):
    cell.text = str(text)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.01)
    cell.margin_right = Inches(0.01)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color

# ============================================================
# ANALISIS AUTOMATICO
# ============================================================

def summarize_values(forecast):
    temps = [valnum(b.get("temperatura"), None) for b in forecast]
    vtos = [valnum(b.get("viento"), None) for b in forecast]
    rachas = [valnum(b.get("rachas"), None) for b in forecast]
    vis = [valnum(b.get("visibilidad"), None) for b in forecast]
    prec = [valnum(b.get("precipitacion"), 0) for b in forecast]
    hums = [get_humedad(b) for b in forecast]
    temps = [x for x in temps if x is not None]
    vtos = [x for x in vtos if x is not None]
    rachas = [x for x in rachas if x is not None]
    vis = [x for x in vis if x is not None]
    hums = [x for x in hums if x is not None]
    return {
        "tmin": min(temps) if temps else None,
        "tmax": max(temps) if temps else None,
        "vmax": max(vtos) if vtos else None,
        "rmax": max(rachas) if rachas else None,
        "vismin": min(vis) if vis else None,
        "precmax": max(prec) if prec else 0,
        "hrmin": min(hums) if hums else None,
        "hrmax": max(hums) if hums else None,
    }


def find_windows(forecast, wanted="FAVORABLE"):
    windows = []
    start = None
    end = None
    cur_date = None
    for b in forecast:
        estado = normalize(estado_general_bloque(b))
        ok = (wanted == "FAVORABLE" and ("FAVORABLE" in estado or estado == "GO"))
        hora = format_dt(b.get("hora"))
        fecha = format_dt(b.get("hora"), "daylabel")
        if ok:
            if start is None:
                start = hora
                cur_date = fecha
            end = hora
        else:
            if start is not None:
                windows.append((cur_date, start, end))
                start = None
                end = None
                cur_date = None
    if start is not None:
        windows.append((cur_date, start, end))
    return windows


def critical_events(data):
    forecast = data.get("forecast72h", []) or []
    events = []
    windows = find_windows(forecast)
    if windows:
        d, ini, fin = max(windows, key=lambda x: int(x[2] or 0) - int(x[1] or 0) if str(x[2]).isdigit() and str(x[1]).isdigit() else 0)
        events.append({"estado": "FAVORABLE", "titulo": "MEJOR VENTANA", "ventana": f"{d} {ini}-{fin}h", "causa": "Mayor continuidad de bloques favorables", "detalle": "Concentrar operaciones principales dentro de esta ventana."})

    for b in forecast:
        vis = valnum(b.get("visibilidad"), None)
        h = get_humedad(b)
        if vis is not None and vis <= 5:
            events.append({"estado": "RESTRINGIDO", "titulo": "VISIBILIDAD REDUCIDA", "ventana": f"{format_dt(b.get('hora'), 'daylabel')} {format_dt(b.get('hora'))}h", "causa": "Visibilidad crítica", "detalle": f"Visibilidad estimada {vis:g} km. Reevaluar operaciones aéreas y navegación."})
            break
        if h is not None and h >= 95:
            events.append({"estado": "MARGINAL", "titulo": "HUMEDAD MUY ALTA", "ventana": f"{format_dt(b.get('hora'), 'daylabel')} {format_dt(b.get('hora'))}h", "causa": "HR elevada", "detalle": f"HR {h:.0f}%. Posible niebla/neblina o baja visibilidad local."})
            break

    max_gust = None
    max_block = None
    for b in forecast:
        g = valnum(b.get("rachas"), None)
        if g is not None and (max_gust is None or g > max_gust):
            max_gust = g
            max_block = b
    if max_gust is not None and max_gust >= 13:
        estado = "RESTRINGIDO" if max_gust >= 17 else "MARGINAL"
        events.append({"estado": estado, "titulo": "RACHAS RELEVANTES", "ventana": f"{format_dt(max_block.get('hora'), 'daylabel')} {format_dt(max_block.get('hora'))}h", "causa": "Viento/rachas", "detalle": f"Racha máxima estimada {max_gust:.1f} kt. Revisar límites por operación."})

    prec_blocks = [b for b in forecast if valnum(b.get("precipitacion"), 0) >= 0.1]
    if prec_blocks:
        b = max(prec_blocks, key=lambda x: valnum(x.get("precipitacion"), 0))
        p = valnum(b.get("precipitacion"), 0)
        events.append({"estado": "RESTRINGIDO" if p >= 1 else "MARGINAL", "titulo": "PRECIPITACIÓN", "ventana": f"{format_dt(b.get('hora'), 'daylabel')} {format_dt(b.get('hora'))}h", "causa": "Precipitación", "detalle": f"Precipitación máxima estimada {p:.1f} mm. Afecta saltos y EVACAM."})

    periodo = data.get("periodo_briefing") or {}
    if periodo.get("mayor_72h"):
        events.insert(0, {"estado": "MARGINAL", "titulo": "PROYECCIÓN PRELIMINAR", "ventana": periodo_label(data, forecast), "causa": "Horizonte >72h", "detalle": periodo.get("advertencia") or "Actualizar diariamente para confirmar cambios."})

    return events[:7]


def global_status(forecast):
    if not forecast:
        return "SIN DATOS"
    estados = [normalize(estado_general_bloque(b)) for b in forecast]
    rest = sum(1 for e in estados if "RESTRINGIDO" in e or "NO GO" in e)
    marg = sum(1 for e in estados if "MARGINAL" in e)
    if rest > len(estados) * 0.35:
        return "RESTRINGIDO"
    if rest or marg > len(estados) * 0.25:
        return "MARGINAL"
    return "FAVORABLE"



def estado_score(estado):
    e = normalize(estado)
    if "RESTRINGIDO" in e or "NO GO" in e or e == "REST":
        return 2
    if "MARGINAL" in e or e == "MARG":
        return 1
    return 0


def block_score(block):
    """Puntaje simple: mayor = peor condición operacional."""
    score = estado_score(estado_general_bloque(block)) * 100
    rachas = valnum(block.get("rachas"), 0) or 0
    viento = valnum(block.get("viento"), 0) or 0
    precip = valnum(block.get("precipitacion"), 0) or 0
    vis = valnum(block.get("visibilidad"), 99) or 99
    nub = valnum(block.get("nubosidad"), 0) or 0
    temp = valnum(block.get("temperatura"), 15) or 15
    if rachas >= 17 or viento >= 17:
        score += 35
    elif rachas >= 13 or viento >= 13:
        score += 18
    if precip >= 1:
        score += 35
    elif precip >= 0.1:
        score += 15
    if vis <= 5:
        score += 30
    elif vis <= 10:
        score += 12
    if nub >= 90:
        score += 8
    if temp >= 25:
        score += 10
    return score


def best_window(forecast):
    """Devuelve mejor ventana continua favorable. Si hay varias, elige la de menor riesgo meteo."""
    candidates = []
    cur = []
    for b in forecast or []:
        if estado_score(estado_general_bloque(b)) == 0:
            cur.append(b)
        else:
            if cur:
                candidates.append(cur)
            cur = []
    if cur:
        candidates.append(cur)
    if not candidates:
        return None
    def cand_key(win):
        avg_bad = sum(block_score(b) for b in win) / max(len(win), 1)
        return (len(win), -avg_bad)
    win = max(candidates, key=cand_key)
    return {
        "inicio": win[0],
        "fin": win[-1],
        "dur": len(win),
        "label": f"{format_dt(win[0].get('hora'), 'daylabel')} {format_dt(win[0].get('hora'))}-{format_dt(win[-1].get('hora'))}h",
    }


def worst_block(forecast):
    """Devuelve peor bloque solo si realmente hay condición marginal/restringida o umbral meteo relevante."""
    if not forecast:
        return None, 0
    worst = max(forecast, key=block_score)
    sc = block_score(worst)
    # si todo es favorable y no hay umbral relevante, no mostrar como peor condición
    if estado_score(estado_general_bloque(worst)) == 0 and sc < 18:
        return None, sc
    return worst, sc


def motivo_operacion(block, opname):
    target = normalize(opname)
    original_motivo = ""
    for op in block.get("operaciones", []) or []:
        if target in normalize(op.get("operacion", "")):
            original_motivo = op.get("motivo", "") or ""
            break
    forced = forced_operation_reason(block, opname)
    if forced:
        if original_motivo and normalize(forced) not in normalize(original_motivo):
            return f"{original_motivo} / {forced}"
        return forced or original_motivo
    return original_motivo


def limiting_codes(block, opname=None):
    """Abreviaciones simples para que el personal no especialista entienda por qué se limita una operación."""
    motivo = normalize(motivo_operacion(block, opname) if opname else "")
    codes = []
    viento = valnum(block.get("viento"), 0) or 0
    rachas = valnum(block.get("rachas"), 0) or 0
    precip = valnum(block.get("precipitacion"), 0) or 0
    vis = valnum(block.get("visibilidad"), None)
    nub = valnum(block.get("nubosidad"), 0) or 0
    temp = valnum(block.get("temperatura"), None)
    techo = valnum(get_metric(block, ["techo", "cloud_base", "base_nubes", "baseNubes"], None), None)
    nieve = valnum(block.get("nieve_cm_h", block.get("snowfall")), 0) or 0
    oleaje = valnum(block.get("oleaje", block.get("altura_ola")), 0) or 0

    if "VIENTO" in motivo or "RACHA" in motivo or viento >= 13 or rachas >= 13:
        codes.append("RV")
    if "PRECIP" in motivo or precip >= 0.1:
        codes.append("P")
    if "VIS" in motivo or "NIEBLA" in motivo or (vis is not None and vis <= 10):
        codes.append("V")
    if "NUB" in motivo or nub >= 90 or (techo is not None and techo > 0 and techo < 1000):
        codes.append("N")
    if temp is not None and temp >= 25:
        codes.append("T")
    if nieve > 0 and (viento >= 15 or rachas >= 20):
        codes.append("VB")
    if oleaje >= 1.5 or "OLEAJE" in motivo or "MAR" in motivo:
        codes.append("O")

    seen = []
    for c in codes:
        if c not in seen:
            seen.append(c)
    return "\n".join(seen[:3])


def plain_reason_from_codes(codes):
    if not codes:
        return "Condición meteorológica más exigente del período."
    text = str(codes).replace("\n", ", ")
    mapping = {
        "RV": "rachas/viento",
        "P": "precipitación",
        "V": "visibilidad reducida",
        "N": "nubosidad alta/baja condición visual",
        "T": "temperatura elevada",
        "VB": "viento blanco/nieve",
        "O": "oleaje/estado de mar",
    }
    parts = [mapping.get(x.strip(), x.strip()) for x in text.split(",") if x.strip()]
    return ", ".join(parts) if parts else "Condición meteorológica relevante."

# ============================================================
# MUNICION / INCENDIO
# ============================================================
MUNICIONES = [
    ("TRAZADORA", "Fósforo incandescente"),
    ("HUMO", "Hexacloroetano"),
    ("GRANADAS", "Fragmentación térmica"),
]


def riesgo_incendio_bloque(b):
    temp = valnum(b.get("temperatura"), None)
    viento = valnum(b.get("viento"), None)
    hr = get_humedad(b)
    criterios = 0
    if temp is not None and temp > 25:
        criterios += 1
    if viento is not None and viento > 15:
        criterios += 1
    if hr is not None and hr < 35:
        criterios += 1

    # Matriz de tiro / riesgo de incendio actualizada:
    # 0-1 criterio = favorable; 2 criterios = marginal; 3 criterios = restringido.
    if criterios >= 3:
        return "RESTRINGIDO", RED, criterios
    if criterios >= 2:
        return "MARGINAL", ORANGE, criterios
    return "FAVORABLE", GREEN, criterios

# ============================================================
# SLIDES
# ============================================================

def _clock_interval_for_block(sorted_blocks, index):
    """Devuelve intervalo 0-24h para que el anillo quede cerrado.
    Cada bloque se toma como representativo del tramo entre puntos medios
    con el bloque anterior y el siguiente. El primero parte en 00h y el
    último termina en 24h, evitando zonas sin color.
    """
    if not sorted_blocks:
        return 0.0, 24.0

    dt = _parse_dt_safe(sorted_blocks[index].get("hora"))
    if not dt:
        return 0.0, 24.0

    h = dt.hour + dt.minute / 60.0

    if len(sorted_blocks) == 1:
        return 0.0, 24.0

    if index == 0:
        start = 0.0
    else:
        prev = _parse_dt_safe(sorted_blocks[index - 1].get("hora"))
        if prev and prev.date() == dt.date():
            hp = prev.hour + prev.minute / 60.0
            start = (hp + h) / 2.0
        else:
            start = 0.0

    if index == len(sorted_blocks) - 1:
        end = 24.0
    else:
        nxt = _parse_dt_safe(sorted_blocks[index + 1].get("hora"))
        if nxt and nxt.date() == dt.date():
            hn = nxt.hour + nxt.minute / 60.0
            end = (h + hn) / 2.0
        else:
            end = 24.0

    start = max(0.0, min(24.0, start))
    end = max(0.0, min(24.0, end))
    if end <= start:
        end = min(24.0, start + 3.0)
    return start, end


def _make_tactical_clock_core_image(data, out_path):
    """Solo genera el reloj circular. Los textos y leyendas se dibujan como objetos PPT editables."""
    forecast = data.get("forecast72h", []) or []
    dias = list(grouped_by_date(forecast).items())
    if not dias or Image is None or ImageDraw is None:
        return None

    # Imagen con fondo transparente para insertarla reducida dentro de portada.
    W, H = 640, 470
    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    white = (245, 248, 252, 255)
    muted = (176, 190, 205, 255)
    blue = (83, 169, 255, 255)
    dark = (12, 20, 31, 255)
    ring_bg = (25, 38, 55, 255)
    f_tiny = _pil_font(18, False)
    f_bold = _pil_font(18, True)
    f_center = _pil_font(30, True)

    cx, cy = 315, 230
    outer = 168
    gap = 7
    dias = dias[:6]
    ring_w = max(16, min(28, int((outer - 45) / max(len(dias), 1)) - gap))

    # marcas principales
    for hour, lab in [(0, "00"), (6, "06"), (12, "12"), (18, "18")]:
        ang = math.radians(-90 + hour / 24 * 360)
        x1 = cx + math.cos(ang) * (outer + 5)
        y1 = cy + math.sin(ang) * (outer + 5)
        x2 = cx + math.cos(ang) * (outer + 21)
        y2 = cy + math.sin(ang) * (outer + 21)
        draw.line((x1, y1, x2, y2), fill=muted, width=2)
        tx = cx + math.cos(ang) * (outer + 38)
        ty = cy + math.sin(ang) * (outer + 38)
        _draw_centered(draw, lab, (tx, ty), f_tiny, white)

    for i, (fecha, bloques) in enumerate(dias):
        r_outer = outer - i * (ring_w + gap)
        r_inner = max(8, r_outer - ring_w)
        draw.ellipse((cx-r_outer, cy-r_outer, cx+r_outer, cy+r_outer), fill=ring_bg, outline=(72,88,110,255), width=1)
        draw.ellipse((cx-r_inner, cy-r_inner, cx+r_inner, cy+r_inner), fill=(0,0,0,0))
        sorted_blocks = sorted(bloques, key=lambda b: str(b.get("hora")))
        for j, b in enumerate(sorted_blocks):
            dt = _parse_dt_safe(b.get("hora"))
            if not dt:
                continue
            h1, h2 = _clock_interval_for_block(sorted_blocks, j)
            col = _clock_status_color(estado_general_bloque(b)) + (255,)
            start_ang = -90 + (h1 / 24) * 360
            end_ang = -90 + (h2 / 24) * 360
            draw.pieslice((cx-r_outer, cy-r_outer, cx+r_outer, cy+r_outer), start=start_ang, end=end_ang, fill=col)
            draw.pieslice((cx-r_inner, cy-r_inner, cx+r_inner, cy+r_inner), start=start_ang-0.5, end=end_ang+0.5, fill=(0,0,0,0))
        draw.ellipse((cx-r_outer, cy-r_outer, cx+r_outer, cy+r_outer), outline=(220,230,240,255), width=1)
        draw.ellipse((cx-r_inner, cy-r_inner, cx+r_inner, cy+r_inner), outline=(18,28,40,255), width=1)
        # día como etiqueta editable no es posible dentro de imagen; se deja pequeño, el resto de texto va en PPT.
        if sorted_blocks:
            label = format_dt(sorted_blocks[0].get("hora"), "daylabel")
            draw.text((470, 110+i*35), label, font=f_bold, fill=white)
            draw.line((440, 120+i*35, 463, 120+i*35), fill=muted, width=1)

    draw.ellipse((cx-42, cy-42, cx+42, cy+42), fill=dark, outline=blue, width=2)
    _draw_centered(draw, "24h", (cx, cy-8), f_center, white)
    _draw_centered(draw, "día", (cx, cy+22), f_tiny, muted)

    # aguja de ahora si cae dentro del período
    now = datetime.now()
    for _, bloques in dias:
        dts = [_parse_dt_safe(b.get("hora")) for b in bloques if _parse_dt_safe(b.get("hora"))]
        if dts and min(dts).date() <= now.date() <= max(dts).date():
            h = now.hour + now.minute / 60
            ang = math.radians(-90 + h/24*360)
            draw.line((cx, cy, cx + math.cos(ang)*(outer+2), cy + math.sin(ang)*(outer+2)), fill=blue, width=4)
            draw.ellipse((cx-5, cy-5, cx+5, cy+5), fill=blue)
            break

    img.save(out_path)
    return out_path


def slide_portada(prs, data):
    forecast = data.get("forecast72h", []) or []
    meteo = data.get("meteo_actual", {}) or {}
    s = summarize_values(forecast)
    status = global_status(forecast)
    events = critical_events(data)
    hours = period_hours_from_forecast(forecast)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_suffix = f"Proyección {hours} horas" if hours else periodo_tipo(data)
    add_title(slide, data, "Slide 1 · Centro de Operaciones Meteorológicas", title_suffix)

    # Alertas críticas
    add_rect(slide, 0.35, 1.55, 12.65, 0.70, PANEL2, RGBColor(50,70,95))
    add_rect(slide, 0.35, 1.55, 0.06, 0.70, RED, RED, False)
    add_text(slide, "⚠  A L E R T A S   C R Í T I C A S   D E L   P E R Í O D O", 0.55, 1.68, 8.0, 0.18, 9.5, True, RED)
    if events:
        lines = [f"• {e['titulo']}: {text_wrap(e['detalle'], 105)}" for e in events[:3]]
    else:
        lines = ["• Sin alertas críticas detectadas para el período evaluado."]
    add_text(slide, "\n".join(lines), 0.55, 1.92, 12.0, 0.28, 7.8, False, WHITE)

    # Mejor / peor ventana
    bw = best_window(forecast)
    wb, _ = worst_block(forecast)
    add_rect(slide, 0.35, 2.40, 6.25, 1.55, PANEL2, GREEN)
    add_text(slide, "✓  M E J O R   V E N T A N A", 0.65, 2.57, 3.4, 0.2, 10.5, True, GREEN)
    if bw:
        add_text(slide, bw["label"], 0.65, 2.98, 5.6, 0.48, 24, True, GREEN)
        add_text(slide, "Concentrar operaciones principales dentro de esta ventana.", 0.65, 3.48, 5.6, 0.22, 8.6, False, WHITE)
        add_text(slide, "Ventana con menor restricción operacional detectada.", 0.65, 3.70, 5.6, 0.18, 7.8, True, YELLOW)
    else:
        add_text(slide, "Sin ventana favorable", 0.65, 3.00, 5.6, 0.4, 20, True, ORANGE)
        add_text(slide, "Planificar con restricciones y actualizar antes de ejecutar.", 0.65, 3.47, 5.6, 0.3, 8.5, False, MUTED)

    add_rect(slide, 6.85, 2.40, 6.15, 1.55, PANEL2, RED)
    add_text(slide, "✗  P E O R   V E N T A N A", 7.15, 2.57, 3.4, 0.2, 10.5, True, RED)
    if wb:
        add_text(slide, f"{format_dt(wb.get('hora'), 'daylabel')} {format_dt(wb.get('hora'))}h", 7.15, 2.98, 5.55, 0.48, 24, True, RED)
        add_text(slide, f"Estado: {estado_general_bloque(wb)}", 7.15, 3.48, 5.55, 0.2, 8.6, False, WHITE)
        add_text(slide, f"Causa: {plain_reason_from_codes(limiting_codes(wb))}", 7.15, 3.70, 5.55, 0.18, 7.8, True, YELLOW)
    else:
        add_text(slide, "Sin condición crítica", 7.15, 3.00, 5.55, 0.4, 20, True, GREEN)
        add_text(slide, "No se detectan restricciones meteorológicas relevantes.", 7.15, 3.47, 5.55, 0.3, 8.5, False, MUTED)

    # Reloj táctico reducido. Solo el reloj es imagen; textos/leyenda son editables.
    add_rect(slide, 0.35, 4.15, 6.25, 2.75, PANEL2, RGBColor(50,70,95))
    add_text(slide, f"⊙  R E L O J   T Á C T I C O   {hours or 72}   H O R A S", 0.65, 4.32, 4.8, 0.20, 9.8, True, CYAN)
    clock_path = BASE_DIR / "_reloj_tactico_core.png"
    created = _make_tactical_clock_core_image(data, clock_path)
    if created and Path(created).exists():
        slide.shapes.add_picture(str(created), Inches(1.55), Inches(4.55), width=Inches(2.7), height=Inches(2.0))
    else:
        add_text(slide, "Sin datos para reloj", 1.75, 5.25, 2.5, 0.3, 12, True, RED)
    add_text(slide, "Cómo leer:", 4.35, 4.62, 1.4, 0.18, 8.5, True, CYAN)
    add_text(slide, "· Cada anillo = 1 día\n· 12 en punto = 00h\n· Avance horario = 00→24h\n· Color = peor estado/ops\n· Aguja azul = AHORA", 4.35, 4.88, 1.95, 0.85, 6.7, False, WHITE)
    legend = [("Óptimo", GREEN2), ("Fav", GREEN), ("Marg", RGBColor(242,204,18)), ("Restr", RED)]
    lx = 0.65
    for label, col in legend:
        add_rect(slide, lx, 6.58, 0.16, 0.16, col, col, False)
        add_text(slide, label, lx+0.28, 6.54, 0.75, 0.18, 6.7, False, WHITE)
        lx += 1.20

    # Tablero actual + luna
    add_rect(slide, 6.85, 4.15, 6.15, 2.75, PANEL2, RGBColor(50,70,95))
    add_text(slide, "⚙  T A B L E R O   D E   I N S T R U M E N T O S  —  E S T A D O   A C T U A L", 7.15, 4.32, 5.55, 0.20, 9.5, True, CYAN)
    vals = [
        ("VIENTO", fmt_num(meteo.get("viento"), 0, ""), "kt"),
        ("RACHAS", fmt_num(meteo.get("rachas"), 0, ""), "kt"),
        ("VISIB", fmt_num(meteo.get("visibilidad"), 0, ""), "km"),
    ]
    for i, (k, v, unit) in enumerate(vals):
        x = 7.15 + i * 1.82
        add_rect(slide, x, 4.78, 1.55, 1.05, DARK, GREEN)
        add_text(slide, k, x+0.10, 4.96, 1.35, 0.18, 7.4, True, MUTED, PP_ALIGN.CENTER)
        add_text(slide, v, x+0.10, 5.25, 1.35, 0.30, 20, True, WHITE, PP_ALIGN.CENTER)
        add_text(slide, unit, x+0.10, 5.58, 1.35, 0.16, 6.5, False, MUTED, PP_ALIGN.CENTER)
        add_rect(slide, x+0.12, 5.74, 1.31, 0.04, GREEN, GREEN, False)

    add_text(slide, "🌙 ILUMINACIÓN LUNAR", 7.15, 6.05, 2.4, 0.18, 9, True, CYAN)
    lunar = get_lunar_source(data)
    illum = lunar.get("iluminacion")
    illum_txt = f"{illum:.0f}%" if isinstance(illum, (int, float)) else "N/D"
    fase = lunar.get("fase", "Fase lunar estimada")
    add_text(slide, illum_txt, 7.15, 6.28, 1.1, 0.28, 16, True, GREEN)
    add_text(slide, text_wrap(fase, 56), 8.0, 6.30, 4.6, 0.20, 7.3, False, GREEN)
    add_text(slide, "Resumen", 11.95, 6.82, 0.7, 0.12, 6.5, False, MUTED)

def slide_heatmap(prs, data):
    forecast = (data.get("forecast72h", []) or [])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, data, "Heatmap meteorológico", "Bloques horarios del período evaluado")

    metrics = [
        ("VIENTO", ["viento"], "kt"),
        ("RACHAS", ["rachas"], "kt"),
        ("TECHO", ["techo", "cloud_base", "base_nubes"], "ft"),
        ("PRECIP", ["precipitacion"], "mm"),
        ("VISIB", ["visibilidad"], "km"),
        ("TEMP", ["temperatura"], "°C"),
        ("HR", ["humedad", "humedad_relativa", "hr"], "%"),
    ]
    rows = len(metrics) + 2
    cols = len(forecast) + 1 if forecast else 2
    table = slide.shapes.add_table(rows, cols, Inches(0.25), Inches(1.55), Inches(12.85), Inches(4.8)).table
    table.columns[0].width = Inches(1.05)
    if cols > 1:
        for c in range(1, cols):
            table.columns[c].width = Inches(11.75 / max(cols-1, 1))

    try:
        table.cell(0,0).merge(table.cell(1,0))
    except:
        pass

    set_cell(table.cell(0, 0), "", 6, True, DARK)

    if forecast:
        c = 1
        while c < cols:
            b = forecast[c-1]
            day = format_dt(b.get('hora'), 'daylabel')
            start_c = c

            while c + 1 < cols and format_dt(forecast[c].get('hora'), 'daylabel') == day:
                c += 1

            end_c = c

            try:
                if end_c > start_c:
                    table.cell(0, start_c).merge(table.cell(0, end_c))
            except:
                pass

            set_cell(table.cell(0, start_c), day, 8.5, True, DARK)

            for cc in range(start_c, end_c + 1):
                bb = forecast[cc-1]
                set_cell(
                    table.cell(1, cc),
                    format_dt(bb.get('hora')),
                    7.5,
                    True,
                    DARK
                )

            c += 1

    for r, (name, keys, unit) in enumerate(metrics, start=2):
        set_cell(table.cell(r, 0), name, 10, True, PANEL3, align=PP_ALIGN.LEFT)
        for c in range(1, cols):
            b = forecast[c-1] if c-1 < len(forecast) else {}
            raw = get_metric(b, keys, "-")
            n = valnum(raw, None)
            label = "-" if n is None else f"{n:.0f}" if unit != "mm" else f"{n:.1f}"
            fill = PANEL
            if n is None:
                fill = PANEL3
            elif name == "VIENTO":
                fill = RED if n >= 17 else ORANGE if n >= 13 else GREEN
            elif name == "RACHAS":
                fill = RED if n >= 17 else ORANGE if n >= 13 else GREEN
            elif name == "PRECIP":
                fill = RED if n >= 1 else ORANGE if n >= 0.1 else GREEN
            elif name == "VISIB":
                fill = RED if n <= 5 else ORANGE if n <= 10 else GREEN
            elif name == "TEMP":
                fill = ORANGE if n >= 25 else GREEN
            elif name == "HR":
                fill = ORANGE if n >= 95 else GREEN
            elif name == "TECHO":
                fill = RED if n < 500 else ORANGE if n < 1000 else GREEN
            set_cell(table.cell(r, c), label, 10.5, True, fill)

    add_text(slide, "Óptimo/Favorable", 0.45, 6.65, 1.5, 0.2, 8, True, GREEN)
    add_text(slide, "Marginal", 2.05, 6.65, 1.0, 0.2, 8, True, ORANGE)
    add_text(slide, "Restringido", 3.1, 6.65, 1.2, 0.2, 8, True, RED)
    add_text(slide, "Unidades: viento/rachas kt · techo ft · precip mm · vis km · temp °C · HR %", 4.4, 6.65, 8.0, 0.2, 8, False, MUTED)



def _rgb_tuple(color):
    try:
        return (int(color[0]), int(color[1]), int(color[2]))
    except Exception:
        return (255, 255, 255)


def _clock_status_color(est):
    s = normalize(est)
    if "RESTRING" in s or "NO GO" in s:
        return (215, 55, 68)
    if "MARG" in s:
        return (239, 145, 38)
    if "OPT" in s:
        return (40, 205, 126)
    if "FAV" in s or "GO" in s:
        return (38, 180, 100)
    return (100, 112, 128)


def _pil_font(size=24, bold=False):
    if ImageFont is None:
        return None
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for c in candidates:
        try:
            return ImageFont.truetype(c, size)
        except Exception:
            pass
    return ImageFont.load_default()


def _parse_dt_safe(value):
    if isinstance(value, datetime):
        return value
    if not value:
        return None
    s = str(value).replace('Z', '')
    for fmt in (None, "%Y-%m-%dT%H:%M:%S", "%Y-%m-%d %H:%M:%S", "%Y-%m-%dT%H:%M", "%Y-%m-%d %H:%M", "%Y-%m-%d"):
        try:
            if fmt is None:
                return datetime.fromisoformat(s)
            return datetime.strptime(s, fmt)
        except Exception:
            continue
    return None


def _draw_centered(draw, text, xy, font, fill=(245,248,252)):
    try:
        bbox = draw.textbbox((0,0), text, font=font)
        w, h = bbox[2]-bbox[0], bbox[3]-bbox[1]
    except Exception:
        w, h = draw.textsize(text, font=font)
    draw.text((xy[0]-w/2, xy[1]-h/2), text, font=font, fill=fill)


def _make_tactical_clock_image(data, out_path):
    """Genera un reloj táctico circular tipo referencia: cada anillo = un día."""
    forecast = data.get("forecast72h", []) or []
    dias = list(grouped_by_date(forecast).items())
    if not dias:
        return None

    if Image is None or ImageDraw is None:
        return None

    # Mantener legible: hasta 6 anillos; si hay más, se resume en los primeros 6 días.
    dias = dias[:6]
    W, H = 1600, 900
    bg = (9, 15, 24)
    panel = (16, 27, 42)
    line = (78, 98, 122)
    white = (245, 248, 252)
    muted = (176, 190, 205)
    blue = (83, 169, 255)

    img = Image.new("RGB", (W, H), bg)
    draw = ImageDraw.Draw(img)
    f_title = _pil_font(34, True)
    f_sub = _pil_font(18, False)
    f_small = _pil_font(18, False)
    f_small_b = _pil_font(18, True)
    f_tiny = _pil_font(15, False)
    f_center = _pil_font(26, True)

    # Títulos
    draw.text((55, 24), "RELOJ TÁCTICO", font=f_title, fill=white)
    draw.text((58, 72), "Cada anillo = 1 día · color = peor condición operacional", font=f_sub, fill=muted)
    draw.line((55, 105, W-55, 105), fill=blue, width=4)

    # Panel de instrucciones izquierda
    draw.rounded_rectangle((55, 145, 390, 625), radius=24, fill=panel, outline=(55,75,100), width=2)
    draw.text((85, 180), "CÓMO LEER", font=f_small_b, fill=white)
    instructions = [
        "· Cada anillo = 1 día",
        "· 12 en punto = 00h",
        "· Avance horario = 00→24h",
        "· Color = estado general",
        "· Texto al costado = fecha",
    ]
    y = 228
    for t in instructions:
        draw.text((90, y), t, font=f_small, fill=muted)
        y += 35

    # Leyenda
    legend = [("ÓPTIMO/FAV", (38,180,100)), ("MARGINAL", (239,145,38)), ("RESTRINGIDO", (215,55,68)), ("SIN EVAL.", (100,112,128))]
    y = 470
    for label, col in legend:
        draw.rounded_rectangle((85, y, 122, y+26), radius=6, fill=col)
        draw.text((138, y-1), label, font=f_small, fill=white)
        y += 35

    # Reloj circular
    cx, cy = 890, 395
    outer = 220
    gap = 10
    ring_w = max(26, min(46, int((outer-70) / max(len(dias), 1)) - gap))

    # Marcas horarias principales
    for hour, lab in [(0,"00"), (6,"06"), (12,"12"), (18,"18")]:
        ang = math.radians(-90 + hour/24*360)
        x1 = cx + math.cos(ang)*(outer+8)
        y1 = cy + math.sin(ang)*(outer+8)
        x2 = cx + math.cos(ang)*(outer+34)
        y2 = cy + math.sin(ang)*(outer+34)
        draw.line((x1,y1,x2,y2), fill=muted, width=3)
        tx = cx + math.cos(ang)*(outer+58)
        ty = cy + math.sin(ang)*(outer+58)
        _draw_centered(draw, lab, (tx,ty), f_small_b, white)

    # círculos guía suaves
    for i, (fecha, bloques) in enumerate(dias):
        r_outer = outer - i*(ring_w+gap)
        r_inner = r_outer - ring_w
        # fondo de anillo
        draw.ellipse((cx-r_outer,cy-r_outer,cx+r_outer,cy+r_outer), fill=(24,34,49), outline=(72,88,110), width=2)
        draw.ellipse((cx-r_inner,cy-r_inner,cx+r_inner,cy+r_inner), fill=bg)

        # Bloques horarios
        sorted_blocks = sorted(bloques, key=lambda b: str(b.get('hora')))
        for j, b in enumerate(sorted_blocks):
            dt = _parse_dt_safe(b.get('hora'))
            if not dt:
                continue
            h1, h2 = _clock_interval_for_block(sorted_blocks, j)
            start_ang = -90 + (h1/24)*360
            end_ang = -90 + (h2/24)*360
            est = estado_general_bloque(b)
            col = _clock_status_color(est)
            # borde del bloque
            draw.pieslice((cx-r_outer,cy-r_outer,cx+r_outer,cy+r_outer), start=start_ang, end=end_ang, fill=col)
            draw.pieslice((cx-r_inner,cy-r_inner,cx+r_inner,cy+r_inner), start=start_ang-1, end=end_ang+1, fill=bg)

        # borde final del anillo
        draw.ellipse((cx-r_outer,cy-r_outer,cx+r_outer,cy+r_outer), outline=(220,230,240), width=2)
        draw.ellipse((cx-r_inner,cy-r_inner,cx+r_inner,cy+r_inner), outline=(30,40,55), width=2)

        # Etiqueta del día a la derecha del anillo
        label = format_dt(sorted_blocks[0].get('hora'), 'daylabel') if sorted_blocks else str(fecha)
        draw.rounded_rectangle((1245, 175+i*58, 1510, 218+i*58), radius=10, fill=panel, outline=(55,75,100), width=1)
        draw.text((1265, 185+i*58), f"Anillo {i+1}: {label}", font=f_small_b, fill=white)

    # centro
    draw.ellipse((cx-60,cy-60,cx+60,cy+60), fill=panel, outline=blue, width=3)
    _draw_centered(draw, "24h", (cx, cy-10), f_center, white)
    _draw_centered(draw, "por día", (cx, cy+22), f_tiny, muted)

    # Aguja de ahora si cae dentro del período
    now = datetime.now()
    for fecha, bloques in dias:
        dts = [_parse_dt_safe(b.get('hora')) for b in bloques if _parse_dt_safe(b.get('hora'))]
        if dts and min(dts).date() <= now.date() <= max(dts).date():
            h = now.hour + now.minute/60
            ang = math.radians(-90 + h/24*360)
            draw.line((cx,cy,cx+math.cos(ang)*(outer+5),cy+math.sin(ang)*(outer+5)), fill=blue, width=6)
            draw.ellipse((cx-8,cy-8,cx+8,cy+8), fill=blue)
            break

    # Nota inferior
    draw.rounded_rectangle((55, 680, 1545, 810), radius=22, fill=panel, outline=(55,75,100), width=2)
    draw.text((90, 705), "LECTURA OPERACIONAL", font=f_small_b, fill=white)
    draw.text((90, 745), "Verde: operar · Amarillo: operar con precaución · Rojo: restringir o reevaluar. Confirmar con fuentes oficiales antes de ejecutar.", font=f_small, fill=muted)

    img.save(out_path)
    return out_path


def slide_reloj_tactico(prs, data):
    """Reloj táctico circular: cada anillo representa un día del período solicitado."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    img_path = BASE_DIR / "_reloj_tactico_circular.png"
    created = _make_tactical_clock_image(data, img_path)
    if created and Path(created).exists():
        # La imagen ya trae título, leyenda y lectura operacional; se inserta a pantalla completa
        # para evitar textos sobrepuestos o recortes visuales.
        slide.shapes.add_picture(str(created), Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
    else:
        add_title(slide, data, "Reloj táctico 72 horas", "Visualización circular: cada anillo representa un día")
        add_text(slide, "No fue posible generar el reloj táctico circular. Revisar datos del forecast.", 0.55, 2.0, 12, 0.5, 14, True, RED)


def slide_gantt(prs, data):
    forecast = (data.get("forecast72h", []) or [])
    incluir_costera = bool(data.get("costera"))
    incluir_montana = bool(data.get("incluir_montana", True)) and bool(data.get("montana"))
    ops = ["SALTO BASICO", "LANZAMIENTO", "SALTO LIBRE", "EVACAM", "MARCHA"]
    labels = ["SBM", "LC", "SLM", "EVACAM", "MARCHA"]
    if incluir_costera:
        ops.append("EMBARC")
        labels.append("EMB")
    if incluir_montana:
        ops.append("MONTAÑA")
        labels.append("MONTAÑA")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, data, "Gantt táctico por operación", "Distribución horaria de estados operacionales y limitantes principales")

    cols = len(forecast) + 1 if forecast else 2
    # Dos filas de encabezado: fecha arriba y hora abajo, para evitar texto apretado.
    rows = len(ops) + 2
    table = slide.shapes.add_table(rows, cols, Inches(0.25), Inches(1.55), Inches(12.85), Inches(4.45)).table
    table.columns[0].width = Inches(1.15)
    for c in range(1, cols):
        table.columns[c].width = Inches(11.6 / max(cols-1, 1))

    # Encabezado OPS unido verticalmente
    try:
        table.cell(0, 0).merge(table.cell(1, 0))
    except Exception:
        pass
    set_cell(table.cell(0, 0), "OPS", 7, True, DARK)

    # Fila 0: fecha agrupada por día; Fila 1: hora
    if forecast:
        c = 1
        while c < cols:
            b = forecast[c-1]
            day = format_dt(b.get('hora'), 'daylabel')
            start_c = c
            while c + 1 < cols and format_dt(forecast[c].get('hora'), 'daylabel') == day:
                c += 1
            end_c = c
            try:
                if end_c > start_c:
                    table.cell(0, start_c).merge(table.cell(0, end_c))
            except Exception:
                pass
            set_cell(table.cell(0, start_c), day, 8.5, True, DARK)
            for cc in range(start_c, end_c + 1):
                bb = forecast[cc-1]
                set_cell(
                table.cell(1, cc),
                format_dt(bb.get('hora')),
                7.8,
                True,
                DARK
            )
            c += 1
    else:
        set_cell(table.cell(0, 1), "SIN DATOS", 6, True, DARK)
        set_cell(table.cell(1, 1), "--", 6, True, DARK)

    # Filas de operaciones
    for r, (op, lab) in enumerate(zip(ops, labels), start=2):
        set_cell(table.cell(r, 0), lab, 10, True, PANEL3, align=PP_ALIGN.LEFT)
        for c in range(1, cols):
            b = forecast[c-1] if c-1 < len(forecast) else {}
            est = get_op_state(b, op)
            codes = limiting_codes(b, op) if estado_score(est) > 0 else ""
            set_cell(table.cell(r, c), codes, 10, True, status_color(est))

    # Leyenda clara para no especialistas
    add_rect(slide, 0.35, 6.10, 3.15, 0.95, PANEL2)
    add_text(slide, "ESTADO", 0.55, 6.22, 1.0, 0.18, 8.3, True, CYAN)
    add_badge(slide, "FAV", 0.55, 6.48, 0.6, 0.22, GREEN)
    add_text(slide, "Favorable", 1.22, 6.49, 0.95, 0.16, 7.3, False, WHITE)
    add_badge(slide, "MARG", 2.15, 6.48, 0.7, 0.22, ORANGE)
    add_badge(slide, "REST", 2.15, 6.73, 0.7, 0.22, RED)
    add_text(slide, "Marginal / Restringido", 2.92, 6.55, 0.55, 0.25, 6.2, False, WHITE)

    add_rect(slide, 3.65, 6.10, 5.55, 0.95, PANEL2)
    add_text(slide, "ABREVIACIONES DE LIMITANTES", 3.85, 6.22, 2.9, 0.18, 8.3, True, CYAN)
    add_text(slide, "RV = rachas/viento   P = precipitación   V = visibilidad\nN = nubosidad/condición visual   T = temperatura   VB = viento blanco/nieve   O = oleaje", 3.85, 6.48, 5.1, 0.35, 7.2, False, WHITE)

    add_rect(slide, 9.35, 6.10, 3.55, 0.95, PANEL2)
    add_text(slide, "LECTURA RÁPIDA", 9.55, 6.22, 1.9, 0.18, 8.3, True, CYAN)
    wb, _ = worst_block(forecast)
    if wb:
        msg = f"Mayor precaución: {format_dt(wb.get('hora'), 'daylabel')} {format_dt(wb.get('hora'))}h · {plain_reason_from_codes(limiting_codes(wb))}."
    else:
        msg = "No se detectan condiciones críticas. Mantener actualización previa a ejecución."
    add_text(slide, msg, 9.55, 6.48, 3.1, 0.35, 7.2, False, WHITE)


def slide_tarjetas(prs, data):
    events = critical_events(data)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, data, "Tarjetas críticas del período", "Eventos ordenados por impacto operacional")
    if not events:
        events = [{"estado": "FAVORABLE", "titulo": "SIN ALERTAS", "ventana": "Periodo completo", "causa": "-", "detalle": "No se identifican eventos críticos automáticos."}]

    positions = [
        (0.35, 1.55), (4.65, 1.55), (8.95, 1.55),
        (0.35, 3.55), (4.65, 3.55), (8.95, 3.55),
    ]
    for i, e in enumerate(events[:6]):
        x, y = positions[i]
        add_rect(slide, x, y, 3.85, 1.65, PANEL2)
        add_badge(slide, estado_abbr(e.get("estado")), x + 0.15, y + 0.15, 0.85, 0.32, status_color(e.get("estado")))
        add_text(slide, e.get("titulo", "EVENTO"), x + 1.1, y + 0.16, 2.5, 0.26, 10, True, WHITE)
        add_text(slide, "VENTANA", x + 0.17, y + 0.56, 0.85, 0.16, 6.5, True, MUTED)
        add_text(slide, e.get("ventana", "-"), x + 1.1, y + 0.52, 2.5, 0.22, 8.5, True, BLUE)
        add_text(slide, "CAUSA", x + 0.17, y + 0.83, 0.85, 0.16, 6.5, True, MUTED)
        add_text(slide, e.get("causa", "-"), x + 1.1, y + 0.80, 2.5, 0.2, 8, False, WHITE)
        add_text(slide, "→ " + text_wrap(e.get("detalle", "-"), 85), x + 0.17, y + 1.12, 3.45, 0.38, 7.7, False, WHITE)

    add_text(slide, f"{len(events)} evento(s) crítico(s) identificado(s).", 0.45, 6.82, 12.0, 0.22, 8.2, False, MUTED)


def slide_costera(prs, data):
    costera = data.get("costera")
    if not costera:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, data, "Módulo costero / Armada", "Información marítima condicional según modo costero")

    fields = [
        ("Zona", costera.get("zona", "-")),
        ("Área", costera.get("area", "-")),
        ("Nubosidad", costera.get("nubosidad", "-")),
        ("Visibilidad", costera.get("visibilidad", "-")),
        ("Viento oficial", costera.get("viento", "-")),
        ("Estado del mar", costera.get("estado_mar", "-")),
        ("Aviso", costera.get("alerta_titulo", "-")),
        ("Área afectada", costera.get("area_afectada", "-")),
        ("Dirección oleaje", costera.get("direccion_oleaje", "-")),
        ("Fuente", costera.get("fuente", "-")),
    ]
    table = slide.shapes.add_table(len(fields)+1, 2, Inches(0.5), Inches(1.55), Inches(12.3), Inches(5.35)).table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(9.8)
    set_cell(table.cell(0,0), "CAMPO", 9, True, DARK)
    set_cell(table.cell(0,1), "DATO / OBSERVACIÓN", 9, True, DARK)
    for r, (k, v) in enumerate(fields, 1):
        fill = PANEL if r % 2 else PANEL2
        set_cell(table.cell(r,0), k, 8.5, True, fill, align=PP_ALIGN.LEFT)
        set_cell(table.cell(r,1), text_wrap(v, 120), 8, False, fill, align=PP_ALIGN.LEFT)




# ============================================================
# MONTAÑA PLUS - HELPERS SOW
# ============================================================
def _montana_estado(block):
    m = block.get("montana") or {}
    return str(m.get("estado") or block.get("estado_montana") or "SIN DATOS")


def _montana_motivo(block):
    m = block.get("montana") or {}
    return str(m.get("motivo") or "Sin datos")


def _montana_score(block):
    estado = normalize(_montana_estado(block))
    score = 0
    if "RESTRINGIDO" in estado:
        score += 100
    elif "MARGINAL" in estado:
        score += 50
    viento = valnum(block.get("viento"), 0) or 0
    rachas = valnum(block.get("rachas"), 0) or 0
    vis = valnum(block.get("visibilidad"), 99) or 99
    nieve = valnum(block.get("nieve_cm_h"), 0) or 0
    acum = valnum(block.get("nieve_acumulada_cm"), 0) or 0
    iso = valnum(block.get("isoterma_0_m"), None)
    techo = valnum(block.get("techo") or block.get("base_nubes_ft"), None)
    temp = valnum(block.get("temperatura"), 15) or 15
    if viento >= 35 or rachas >= 45:
        score += 35
    elif viento >= 20 or rachas >= 30:
        score += 15
    if vis < 2:
        score += 35
    elif vis < 5:
        score += 15
    if nieve >= 2:
        score += 30
    elif nieve >= 0.1:
        score += 10
    if acum >= 15:
        score += 25
    if iso is not None and iso > 0 and iso < 2200:
        score += 20
    if techo is not None and techo > 0 and techo < 1000:
        score += 20
    if temp <= -8:
        score += 20
    return score


def _montana_best_window(forecast):
    if not forecast:
        return "Sin datos"
    best = []
    cur = []
    for b in forecast:
        estado = normalize(_montana_estado(b))
        if "FAVORABLE" in estado:
            cur.append(b)
        else:
            if len(cur) > len(best):
                best = list(cur)
            cur = []
    if len(cur) > len(best):
        best = list(cur)
    if not best:
        return "Sin ventana favorable"
    return f"{format_dt(best[0].get('hora'), 'daylabel')} {format_dt(best[0].get('hora'))}-{format_dt(best[-1].get('hora'))}h"


def _montana_riesgos_from_forecast(forecast):
    if not forecast:
        return []
    def nums(key):
        return [valnum(b.get(key), None) for b in forecast if valnum(b.get(key), None) is not None]
    viento = max(nums("viento") or [0])
    rachas = max(nums("rachas") or [0])
    vis = min(nums("visibilidad") or [0])
    nieve = max(nums("nieve_cm_h") or [0])
    acum = max(nums("nieve_acumulada_cm") or [0])
    iso_vals = [v for v in nums("isoterma_0_m") if v > 0]
    techo_vals = [v for v in (nums("techo") + nums("base_nubes_ft")) if v > 0]
    temp = min(nums("temperatura") or [0])

    def ev_wind():
        if viento >= 35 or rachas >= 45: return "RESTRINGIDO"
        if viento >= 20 or rachas >= 30: return "MARGINAL"
        return "FAVORABLE"
    def ev_vis():
        if vis <= 0: return "SIN DATO"
        if vis < 2: return "RESTRINGIDO"
        if vis < 5: return "MARGINAL"
        return "FAVORABLE"
    def ev_snow():
        if nieve >= 2: return "RESTRINGIDO"
        if nieve >= 0.1: return "MARGINAL"
        return "FAVORABLE"
    def ev_acum():
        if acum >= 15: return "RESTRINGIDO"
        if acum >= 5: return "MARGINAL"
        return "FAVORABLE"
    def ev_iso():
        if not iso_vals: return "SIN DATO"
        m = min(iso_vals)
        if m < 2200: return "RESTRINGIDO"
        if m <= 3000: return "MARGINAL"
        return "FAVORABLE"
    def ev_techo():
        if not techo_vals: return "SIN DATO"
        m = min(techo_vals)
        if m < 1000: return "RESTRINGIDO"
        if m < 2000: return "MARGINAL"
        return "FAVORABLE"
    def ev_temp():
        if temp <= -8: return "RESTRINGIDO"
        if temp <= 0: return "MARGINAL"
        return "FAVORABLE"

    return [
        {"riesgo": "Viento/rachas", "estado": ev_wind(), "detalle": f"{viento:.0f}/{rachas:.0f} kt"},
        {"riesgo": "Visibilidad", "estado": ev_vis(), "detalle": f"{vis:.1f} km mín" if vis else "N/D"},
        {"riesgo": "Nevada", "estado": ev_snow(), "detalle": f"{nieve:.1f} cm/h máx"},
        {"riesgo": "Nieve acumulada", "estado": ev_acum(), "detalle": f"{acum:.1f} cm máx"},
        {"riesgo": "Isoterma 0", "estado": ev_iso(), "detalle": f"{min(iso_vals):.0f} m mín" if iso_vals else "N/D"},
        {"riesgo": "Techo/base nubes", "estado": ev_techo(), "detalle": f"{min(techo_vals):.0f} ft mín" if techo_vals else "N/D"},
        {"riesgo": "Frío/exposición", "estado": ev_temp(), "detalle": f"{temp:.1f}°C mín"},
    ]


def _montana_dias_from_forecast(forecast):
    dias = grouped_by_date(forecast)
    out = []
    for fecha, blocks in dias.items():
        worst = max(blocks, key=_montana_score) if blocks else None
        estados = [normalize(_montana_estado(b)) for b in blocks]
        estado = "RESTRINGIDO" if any("RESTRINGIDO" in e for e in estados) else "MARGINAL" if any("MARGINAL" in e for e in estados) else "FAVORABLE"
        out.append({
            "fecha": fecha,
            "estado": estado,
            "mejor_ventana": _montana_best_window(blocks),
            "peor_bloque": f"{format_dt(worst.get('hora'))}h" if worst else "-",
            "motivo": _montana_motivo(worst) if worst else "Sin datos",
        })
    return out


def _montana_resumen_auto(data):
    forecast = data.get("forecast72h", []) or []
    montana = data.get("montana") or {}
    if not forecast:
        return {
            "estado": montana.get("estado", "SIN DATOS"),
            "motivo": montana.get("motivo", "Sin datos"),
            "mejor_ventana": "Sin datos",
            "peor_bloque": "Sin datos",
            "peor_estado": "SIN DATOS",
            "riesgos": [],
            "dias": [],
            "conclusion": "Sin datos suficientes para evaluar montaña.",
            "fuente": "Open-Meteo / SOW",
        }
    worst = max(forecast, key=_montana_score)
    estados = [normalize(_montana_estado(b)) for b in forecast]
    rest = sum(1 for e in estados if "RESTRINGIDO" in e)
    marg = sum(1 for e in estados if "MARGINAL" in e)
    estado = "RESTRINGIDO" if rest else "MARGINAL" if marg else "FAVORABLE"
    motivo = _montana_motivo(worst)
    if estado == "RESTRINGIDO":
        conclusion = f"Operación de montaña NO recomendable en los bloques restringidos. Motivo principal: {motivo}. Reforzar navegación, abrigo, visibilidad, ruta de evacuación y actualización meteorológica."
    elif estado == "MARGINAL":
        conclusion = f"Operación de montaña posible solo con control permanente. Priorizar mejor ventana y reevaluar viento, visibilidad, nieve e isoterma 0. Motivo principal: {motivo}."
    else:
        conclusion = "Condiciones generales favorables para montaña. Mantener monitoreo por cambios rápidos en altura."
    return {
        "estado": estado,
        "motivo": motivo,
        "bloques_restringidos": rest,
        "bloques_marginales": marg,
        "mejor_ventana": _montana_best_window(forecast),
        "peor_bloque": f"{format_dt(worst.get('hora'), 'daylabel')} {format_dt(worst.get('hora'))}h",
        "peor_estado": _montana_estado(worst),
        "peor_motivo": motivo,
        "riesgos": _montana_riesgos_from_forecast(forecast),
        "dias": _montana_dias_from_forecast(forecast),
        "conclusion": conclusion,
        "fuente": "Open-Meteo / SOW",
    }

def slide_montana(prs, data):
    if not data.get("incluir_montana", True):
        return

    forecast = data.get("forecast72h", []) or []
    resumen = data.get("montana_resumen") or _montana_resumen_auto(data)
    montana = data.get("montana") or {}

    if not resumen and not montana and not forecast:
        return

    estado = str(resumen.get("estado") or montana.get("estado") or "SIN EVALUAR")
    motivo = str(resumen.get("motivo") or montana.get("motivo") or "Sin datos")
    riesgos = resumen.get("riesgos") or _montana_riesgos_from_forecast(forecast)
    dias = resumen.get("dias") or _montana_dias_from_forecast(forecast)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, data, "MODO MONTAÑA Y NIEVE", "Evaluación operacional de isoterma 0, nieve, visibilidad, viento, techo y exposición.")

    # Panel principal
    add_rect(slide, 0.45, 1.55, 3.25, 1.45, PANEL2)
    add_text(slide, "ESTADO MONTAÑA", 0.62, 1.72, 2.7, 0.22, 9.5, True, MUTED)
    add_badge(slide, estado, 0.62, 2.03, 2.65, 0.48, status_color(estado))
    add_text(slide, text_wrap(motivo, 150), 0.62, 2.58, 2.75, 0.32, 7.8, False, WHITE)

    add_rect(slide, 3.95, 1.55, 3.05, 1.45, PANEL)
    add_text(slide, "MEJOR VENTANA", 4.12, 1.72, 2.6, 0.22, 9.5, True, GREEN)
    add_text(slide, str(resumen.get("mejor_ventana", "Sin ventana")), 4.12, 2.05, 2.5, 0.35, 15, True, WHITE)
    add_text(slide, "Priorizar ejecución en este bloque si el resto de factores se mantiene.", 4.12, 2.52, 2.55, 0.33, 7.3, False, MUTED)

    add_rect(slide, 7.22, 1.55, 2.75, 1.45, PANEL)
    add_text(slide, "PEOR BLOQUE", 7.38, 1.72, 2.35, 0.22, 9.5, True, RED)
    add_text(slide, str(resumen.get("peor_bloque", "Sin bloque crítico")), 7.38, 2.05, 2.25, 0.35, 13, True, WHITE)
    add_text(slide, text_wrap(str(resumen.get("peor_motivo", "Sin datos")), 90), 7.38, 2.52, 2.25, 0.33, 7.2, False, MUTED)

    add_rect(slide, 10.18, 1.55, 2.7, 1.45, PANEL)
    add_text(slide, "FUENTE", 10.35, 1.72, 2.2, 0.22, 9.5, True, CYAN)
    add_text(slide, str(resumen.get("fuente", "Open-Meteo / SOW")), 10.35, 2.05, 2.2, 0.3, 12, True, WHITE)
    add_text(slide, "Snow/Mountain-Forecast: referencia manual hasta contar con API autorizada.", 10.35, 2.43, 2.2, 0.42, 6.8, False, MUTED)

    # Datos críticos
    add_rect(slide, 0.45, 3.25, 5.9, 1.45, PANEL)
    add_text(slide, "DATOS CRÍTICOS", 0.62, 3.42, 5.3, 0.22, 10, True, CYAN)
    datos = [
        ("Temp", resumen.get("temp_rango", "-")),
        ("Viento/Rachas", resumen.get("viento_rachas", "-")),
        ("Vis mín", resumen.get("visibilidad_min", "-")),
        ("Nieve máx", resumen.get("nieve_max", "-")),
        ("Nieve acum", resumen.get("nieve_acum_max", "-")),
        ("Isoterma", resumen.get("isoterma_rango", "-")),
        ("Techo mín", resumen.get("techo_min", "-")),
    ]
    for i, (k, v) in enumerate(datos[:7]):
        x = 0.62 + (i % 4) * 1.42
        y = 3.78 + (i // 4) * 0.48
        add_text(slide, str(k).upper(), x, y, 1.25, 0.16, 6.6, True, MUTED)
        add_text(slide, str(v), x, y + 0.16, 1.30, 0.22, 8.2, True, WHITE)

    # Riesgos
    add_rect(slide, 6.6, 3.25, 6.28, 1.45, PANEL)
    add_text(slide, "RIESGOS OPERACIONALES", 6.78, 3.42, 5.6, 0.22, 10, True, CYAN)
    for i, r in enumerate((riesgos or [])[:4]):
        x = 6.78 + (i % 2) * 3.0
        y = 3.78 + (i // 2) * 0.48
        est = str(r.get("estado", "SIN DATO"))
        add_text(slide, str(r.get("riesgo", "-")), x, y, 1.6, 0.18, 7.2, True, WHITE)
        add_badge(slide, estado_abbr(est), x + 1.62, y - 0.01, 0.68, 0.22, status_color(est))
        add_text(slide, str(r.get("detalle", "-")), x + 2.35, y, 0.55, 0.18, 6.5, False, MUTED)

    # Resumen por día
    add_rect(slide, 0.45, 4.95, 12.43, 1.25, PANEL2)
    add_text(slide, "SÍNTESIS POR DÍA", 0.62, 5.10, 3.0, 0.22, 10, True, CYAN)
    x = 0.62
    for d in (dias or [])[:4]:
        est = str(d.get("estado", "SIN DATOS"))
        add_rect(slide, x, 5.42, 2.95, 0.55, PANEL)
        add_text(slide, str(d.get("fecha", "-")), x + 0.08, 5.48, 0.72, 0.16, 6.6, True, MUTED)
        add_badge(slide, estado_abbr(est), x + 0.78, 5.47, 0.70, 0.20, status_color(est))
        add_text(slide, f"Mejor: {d.get('mejor_ventana', '-')}", x + 1.58, 5.47, 1.25, 0.16, 6.4, False, WHITE)
        add_text(slide, text_wrap(f"Peor {d.get('peor_bloque', '-')}: {d.get('motivo', '-')}", 70), x + 0.08, 5.70, 2.72, 0.16, 6.2, False, MUTED)
        x += 3.05

    add_rect(slide, 0.45, 6.35, 12.43, 0.62, DARK)
    add_text(slide, "CONCLUSIÓN", 0.62, 6.43, 1.6, 0.18, 8.8, True, CYAN)
    add_text(slide, text_wrap(str(resumen.get("conclusion", "Mantener monitoreo meteorológico en montaña.")), 270), 1.75, 6.40, 10.75, 0.33, 8.1, False, WHITE)


def slide_municion(prs, data):
    forecast = (data.get("forecast72h", []) or [])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, data, "Matriz munición / riesgo de incendio", "Temperatura, viento y humedad relativa")

    total_rest = 0
    total_marg = 0

    for b in forecast:
        estado, _, _ = riesgo_incendio_bloque(b)
        if estado == "RESTRINGIDO":
            total_rest += 1
        elif estado == "MARGINAL":
            total_marg += 1

    overall = "RESTRINGIDO" if total_rest else "MARGINAL" if total_marg else "SIN RIESGO"

    add_rect(slide, 0.45, 1.55, 4.0, 1.0, PANEL2)
    add_text(slide, "🔥 ÍNDICE DE RIESGO DE INCENDIO", 0.65, 1.73, 3.5, 0.2, 10, True, CYAN)
    add_badge(slide, overall, 0.75, 2.05, 2.6, 0.32, status_color(overall) if overall != "SIN RIESGO" else GREEN)

    add_text(slide, "Criterios: T>25°C · viento>15kt · HR<35%", 4.8, 1.77, 7.5, 0.22, 9.2, False, MUTED)
    add_text(slide, "3 criterios = restringido | 2 criterios = marginal | 0-1 = favorable", 4.8, 2.08, 7.5, 0.22, 9.2, False, MUTED)

    cols = len(forecast) + 1 if forecast else 2
    rows = len(MUNICIONES) + 2

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.35),
        Inches(3.05),
        Inches(12.7),
        Inches(2.35)
    ).table

    table.columns[0].width = Inches(2.3)

    for c in range(1, cols):
        table.columns[c].width = Inches(10.3 / max(cols - 1, 1))

    set_cell(table.cell(0, 0), "TIPO", 8, True, DARK)

    if forecast:
        c = 1

        while c < cols:
            b = forecast[c - 1]
            day = format_dt(b.get("hora"), "daylabel")

            start_c = c

            while c + 1 < cols and format_dt(
                forecast[c].get("hora"),
                "daylabel"
            ) == day:
                c += 1

            end_c = c

            if end_c > start_c:
                try:
                    table.cell(0, start_c).merge(table.cell(0, end_c))
                except Exception:
                    pass

            set_cell(
                table.cell(0, start_c),
                day,
                8,
                True,
                DARK
            )

            for cc in range(start_c, end_c + 1):
                bb = forecast[cc - 1]

                set_cell(
                    table.cell(1, cc),
                    format_dt(bb.get("hora")),
                    7,
                    True,
                    DARK
                )

            c += 1

    for r, (tipo, desc) in enumerate(MUNICIONES, start=2):
        set_cell(
            table.cell(r, 0),
            f"{tipo}\\n{desc}",
            7,
            True,
            PANEL3,
            align=PP_ALIGN.LEFT
        )

        for c in range(1, cols):
            b = forecast[c - 1] if c - 1 < len(forecast) else {}

            estado, color, _ = riesgo_incendio_bloque(b)

            set_cell(
                table.cell(r, c),
                estado_abbr(estado),
                7,
                True,
                color
            )



def _blocks_sorted(blocks):
    return sorted(blocks or [], key=lambda b: str(b.get("hora", "")))


def _fmt_hour(block):
    return format_dt(block.get("hora"))


def _compress_windows(blocks):
    """Convierte bloques horarios consecutivos en ventanas legibles: 02-08h, 14h."""
    blocks = _blocks_sorted(blocks)
    if not blocks:
        return ""

    windows = []
    current = []

    for b in blocks:
        if not current:
            current = [b]
            continue

        prev_dt = _parse_dt_safe(current[-1].get("hora"))
        cur_dt = _parse_dt_safe(b.get("hora"))

        contiguous = True
        if prev_dt and cur_dt:
            diff_h = (cur_dt - prev_dt).total_seconds() / 3600.0
            contiguous = 0 < diff_h <= 3.2

        if contiguous:
            current.append(b)
        else:
            windows.append(current)
            current = [b]

    if current:
        windows.append(current)

    labels = []
    for w in windows:
        h1 = _fmt_hour(w[0])
        h2 = _fmt_hour(w[-1])
        labels.append(f"{h1}h" if h1 == h2 else f"{h1}-{h2}h")

    if len(labels) > 4:
        return ", ".join(labels[:4]) + "..."
    return ", ".join(labels)


def _window_blocks(blocks, predicate):
    out = []
    for b in _blocks_sorted(blocks):
        try:
            if predicate(b):
                out.append(b)
        except Exception:
            pass
    return out


def _op_window(blocks, opname, wanted_rank):
    return _window_blocks(
        blocks,
        lambda b: _state_rank(get_op_state(b, opname)) == wanted_rank
    )


def _op_restricted_window(blocks, opname):
    return _op_window(blocks, opname, 2)


def _op_marginal_window(blocks, opname):
    return _op_window(blocks, opname, 1)


def _op_favorable_window(blocks, opname):
    return _op_window(blocks, opname, 0)


def _all_ops_window(blocks, ops, wanted_rank):
    return _window_blocks(
        blocks,
        lambda b: all(_state_rank(get_op_state(b, op)) == wanted_rank for op in ops)
    )


def _reason_short(blocks, opname=None, max_items=2):
    """Motivo operacional compacto, priorizando el motivo forzado por seguridad."""
    reasons = []

    for b in _blocks_sorted(blocks):
        r = ""
        if opname:
            r = forced_operation_reason(b, opname) or motivo_operacion(b, opname)
        if not r:
            r = plain_reason_from_codes(limiting_codes(b, opname))
        for part in str(r).replace(" / ", "|").split("|"):
            p = part.strip()
            if p and p not in reasons and p != "Condición meteorológica más exigente del período.":
                reasons.append(p)

    if not reasons:
        return "condición meteorológica limitante"

    cleaned = []
    for r in reasons[:max_items]:
        r = r.replace("Viento/rachas sobre límite", "rachas sobre límite")
        r = r.replace("Techo/base de nubes", "base de nubes")
        r = r.replace("Precipitación marginal", "precipitación marginal")
        r = r.replace("Visibilidad crítica", "visibilidad crítica")
        r = r.replace("Visibilidad reducida", "visibilidad reducida")
        cleaned.append(r)

    return " y ".join(cleaned)


def _day_state_from_blocks(blocks):
    estado = "FAVORABLE"
    for b in blocks:
        est = normalize(estado_general_bloque(b))
        if "RESTRINGIDO" in est or "NO GO" in est:
            return "RESTRINGIDO"
        if "MARGINAL" in est:
            estado = "MARGINAL"
    return estado


def _blocks_with_condition(blocks, predicate):
    return _window_blocks(blocks, predicate)


def _dominant_fire_risk(blocks):
    counts = {"FAVORABLE": 0, "MARGINAL": 0, "RESTRINGIDO": 0}
    worst_criterios = 0
    for b in blocks:
        estado, _, criterios = riesgo_incendio_bloque(b)
        counts[estado] = counts.get(estado, 0) + 1
        worst_criterios = max(worst_criterios, criterios)
    if counts["RESTRINGIDO"]:
        return "RESTRINGIDO", worst_criterios
    if counts["MARGINAL"]:
        return "MARGINAL", worst_criterios
    return "FAVORABLE", worst_criterios


def _daily_operational_conclusion(fecha_txt, bloques):
    bloques_validos = [b for b in _blocks_sorted(bloques) if not b.get("sin_datos")]

    if not bloques_validos:
        return (
            "SIN DATOS",
            f"{fecha_txt} — SIN DATOS\n"
            "No se recibieron bloques meteorológicos para este día desde la app/API.\n"
            "Mantener el día solo como referencia de planificación.\n"
            "Actualizar la proyección antes de ordenar la ejecución."
        )

    estado = _day_state_from_blocks(bloques_validos)

    aerial_ops = [
        ("Salto basico militar", "SBM"),
        ("Lanzamiento de carga", "LC"),
        ("Salto libre militar", "SLM"),
        ("EVACAM", "EVACAM"),
    ]

    lines = [f"{fecha_txt} — {estado}"]

    # 1) Decisión operacional aérea y ventanas.
    all_aerial_rest = _all_ops_window(
        bloques_validos,
        [op for op, _ in aerial_ops],
        2
    )
    all_aerial_marg = _all_ops_window(
        bloques_validos,
        [op for op, _ in aerial_ops],
        1
    )
    all_aerial_fav = _all_ops_window(
        bloques_validos,
        [op for op, _ in aerial_ops],
        0
    )

    if all_aerial_rest:
        lines.append(
            f"Aéreas/evacuación: NO GO en {_compress_windows(all_aerial_rest)} por "
            f"{_reason_short(all_aerial_rest, 'Salto basico militar')}; afecta despegue/aterrizaje, navegación y seguridad de lanzamiento."
        )
    else:
        restricciones = []
        for op, abbr in aerial_ops:
            rest = _op_restricted_window(bloques_validos, op)
            if rest:
                restricciones.append(
                    f"{abbr} NO GO {_compress_windows(rest)} ({_reason_short(rest, op)})"
                )
        if restricciones:
            lines.append("Restricciones por operación: " + "; ".join(restricciones[:3]) + ".")
        else:
            lines.append("Aéreas/evacuación: sin bloques NO GO comunes; mantener control de rachas, visibilidad y techo antes de ejecutar.")

    if all_aerial_marg:
        lines.append(
            f"Bloques marginales comunes: {_compress_windows(all_aerial_marg)}; ejecutar solo con decisión del mando, control horario y alternativa terrestre."
        )

    if all_aerial_fav:
        lines.append(
            f"Ventanas recomendadas: {_compress_windows(all_aerial_fav)} para concentrar SBM/LC/SLM/EVACAM y reducir exposición al cambio meteorológico."
        )
    else:
        sbm_fav = _op_favorable_window(bloques_validos, "Salto basico militar")
        evacam_fav = _op_favorable_window(bloques_validos, "EVACAM")
        fav_parts = []
        if sbm_fav:
            fav_parts.append(f"SBM {_compress_windows(sbm_fav)}")
        if evacam_fav:
            fav_parts.append(f"EVACAM {_compress_windows(evacam_fav)}")
        if fav_parts:
            lines.append("Ventanas utilizables parciales: " + "; ".join(fav_parts) + ".")
        else:
            lines.append("No se aprecia ventana aérea favorable continua; planificar ejecución terrestre o reprogramar.")

    # 2) Impacto en tiro, munición y personal.
    temps = [valnum(x.get("temperatura"), None) for x in bloques_validos]
    rachas = [valnum(x.get("rachas"), None) for x in bloques_validos]
    vientos = [valnum(x.get("viento"), None) for x in bloques_validos]
    visibs = [valnum(x.get("visibilidad"), None) for x in bloques_validos]
    precs = [valnum(x.get("precipitacion"), 0) for x in bloques_validos]
    hrs = [get_humedad(x) for x in bloques_validos]
    techos = [valnum(get_metric(x, ["techo", "cloud_base", "base_nubes", "base_nubes_ft", "baseNubes"], None), None) for x in bloques_validos]

    temps = [x for x in temps if x is not None]
    rachas = [x for x in rachas if x is not None]
    vientos = [x for x in vientos if x is not None]
    visibs = [x for x in visibs if x is not None]
    precs = [x for x in precs if x is not None]
    hrs = [x for x in hrs if x is not None]
    techos = [x for x in techos if x is not None and x > 0]

    impactos = []
    if rachas and max(rachas) >= 13:
        impactos.append("rachas relevantes: incrementan riesgo para paracaidistas y obligan a corrección/reevaluación en tiro")
    if visibs and min(visibs) <= 5:
        impactos.append("visibilidad crítica: degrada observación, adquisición de blancos y control de zona")
    elif visibs and min(visibs) <= 10:
        impactos.append("visibilidad reducida: limita observación y seguridad de aproximación")
    if techos and min(techos) < 1000:
        impactos.append("techo bajo: limita apoyo aéreo, navegación y selección de zona de caída")
    if precs and max(precs) >= 0.1:
        impactos.append("precipitación: afecta vialidad, comodidad del personal, óptica y almacenamiento de abastecimientos")
    if temps:
        if min(temps) <= 3:
            impactos.append("frío/baja T°: aumenta exigencia fisiológica, afecta lubricantes, baterías, calentamiento de motores y manipulación de armamento")
        elif max(temps) >= 25:
            impactos.append("calor: aumenta carga fisiológica con equipo completo, demanda hídrica y fatiga")
    if hrs:
        if max(hrs) >= 90:
            impactos.append("HR alta: favorece niebla/neblina y degradación de ópticas/equipos")
        elif min(hrs) < 35:
            impactos.append("HR baja: aumenta sensibilidad de combustible fino para tiro/munición con riesgo de incendio")

    if impactos:
        lines.append("Impacto táctico: " + "; ".join(impactos[:3]) + ".")

    fire_state, fire_crit = _dominant_fire_risk(bloques_validos)
    if fire_state == "RESTRINGIDO":
        lines.append("Matriz tiro/incendio: RESTRINGIDO cuando coinciden T>25°C, viento>15 kt y HR<35%; evitar trazadora/humo/pirotecnia.")
    elif fire_state == "MARGINAL":
        lines.append("Matriz tiro/incendio: MARGINAL por combinación de 2 criterios; autorizar solo con cortafuego, observador y control de viento/HR.")
    else:
        lines.append("Tiro/munición: sin combinación crítica de incendio; mantener control de viento y HR antes de abrir fuego.")

    # 3) Recomendación de mando.
    if estado == "RESTRINGIDO":
        lines.append("Decisión sugerida: suspender bloques NO GO, emplear ventanas favorables puntuales y exigir actualización previa a la orden de ejecución.")
    elif estado == "MARGINAL":
        lines.append("Decisión sugerida: ejecutar solo con mitigaciones, ventanas horarias definidas, plan alterno y monitoreo antes/durante la actividad.")
    else:
        lines.append("Decisión sugerida: día apto; concentrar esfuerzo en ventanas continuas y mantener vigilancia de rachas/visibilidad por cambios locales.")

    # Compactar: no eliminar impacto de mando, pero evitar desbordar excesivamente.
    if len(lines) > 9:
        lines = lines[:8] + [lines[-1]]

    return estado, "\n".join(lines)


def slide_conclusiones(prs, data):
    forecast = data.get("forecast72h", []) or []
    dias = list(grouped_by_date(forecast).items())

    if not dias:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_title(slide, data, "Conclusiones y recomendaciones tácticas", "Síntesis diaria para planificación operacional")
        add_text(slide, "Sin datos de pronóstico disponibles.", 0.55, 1.8, 12, 0.5, 14, True, RED)
        return

    grupos = [dias[i:i+4] for i in range(0, len(dias), 4)]

    for idx, grupo in enumerate(grupos):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)

        subt = "Síntesis operacional por día"
        if len(grupos) > 1:
            subt += f" · Continuación {idx+1}/{len(grupos)}"

        add_title(slide, data, "Conclusiones y recomendaciones tácticas", subt)

        cantidad = len(grupo)

        if cantidad == 1:
            xs = [1.0]
            w = 11.2
        elif cantidad == 2:
            xs = [0.55, 6.92]
            w = 5.85
        elif cantidad == 3:
            xs = [0.35, 4.62, 8.89]
            w = 3.85
        else:
            xs = [0.15, 3.38, 6.61, 9.84]
            w = 3.05

        for i, (fecha, bloques) in enumerate(grupo):
            x = xs[i]
            fecha_txt = str(fecha)

            add_rect(slide, x, 1.55, w, 5.25, PANEL2)

            estado, txt = _daily_operational_conclusion(fecha_txt, bloques)
            lines = txt.split("\n")

            titulo_size = 10
            body_size = 8.6

            color_estado = (
                RED if estado == "RESTRINGIDO"
                else ORANGE if estado == "MARGINAL"
                else MUTED if estado == "SIN DATOS"
                else CYAN
            )

            add_text(slide, lines[0], x+0.12, 1.78, w-0.25, 0.35, titulo_size, True, color_estado)
            add_text(slide, "\n".join(lines[1:]), x+0.12, 2.10, w-0.25, 4.55, body_size, False, WHITE)

        add_text(
            slide,
            "TODAS LAS TEMPERATURAS REGISTRADAS EN LA PROYECCIÓN METEOROLÓGICA SON TOMADAS A LA SOMBRA.",
            0.45,
            6.8,
            12.0,
            0.2,
            8.5,
            True,
            MUTED
        )


# ============================================================
# MAIN
# ============================================================

def main():
    data = apply_period_filter(load_json())
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_portada(prs, data)
    slide_heatmap(prs, data)
    slide_gantt(prs, data)
    slide_tarjetas(prs, data)
    slide_costera(prs, data)
    slide_montana(prs, data)
    slide_municion(prs, data)
    slide_conclusiones(prs, data)

    output_file = output_file_for_data(data)
    prs.save(output_file)
    print(f"PPT generado correctamente: {output_file}")


if __name__ == "__main__":
    main()

